#!/usr/bin/env python3
"""Build a PowerPoint deck out of an "AVANCE BRANCA" (Fratelli Branca) workbook.

Same idea as `scripts/avance_pptx.py`, different sheet shape. In the BADIE
workbook the rows are vendors and the columns are categories; here it is the
other way round: **rows are categorias/marcas and columns are vendors**, three
per vendor (Avance, %Tend, Faltan). The deck keeps the sheet's orientation
because that is how Nahuel reads it. The look and feel is imported from
`avance_pptx` so both decks stay identical on screen.

Deck layout
-----------
1. Cover.
2. "VOLUMEN — LINEA BRANCA": sheet `AVANCE`, rows ANTICA…GRAPPA plus the total.
3. "VOLUMEN — OTRAS LINEAS": sheet `AVANCE`, ARIZU…QUARA.
4. "COBERTURA — LINEA BRANCA": sheet `Cobertura`.
5. "COBERTURA — OTRAS LINEAS": sheet `Cobertura`.
6. "RECHAZOS": the month's rebotes PNG, as an image.

Volume first, coverage after — same reason as the BADIE deck.

Two things this script has to get right
---------------------------------------
**The workbook usually has no cached formula values.** `AvancesService` writes
it with openpyxl, which drops every cached result, so `data_only=True` reads a
sheet full of `None` unless a human opened and saved it since. When that
happens the script recalculates a COPY with LibreOffice (never the original)
and reads that.

**Coverage totals are READ, volume totals are COMPUTED.** Volume adds up across
categories, so `TOTAL LINEA` is the sum of the rows above it. Coverage does NOT:
it counts points of sale, and a client that buys FERNET and CARPANO is one PDV,
not two. In JULIO 2026 the sheet's own `TOTAL LINEA` for PABLO NAVARRO is 201
while its fourteen brand rows add up to 352 — summing there would inflate the
number by 75%. So the coverage total comes from the sheet, verbatim.

Usage
-----
    python scripts/avance_branca_pptx.py \
        --archivo "data/output/avances/2026-07/AVANCE BRANCA - JULIO 2026.xlsx"
"""
from __future__ import annotations

import argparse
import sys
import tempfile
from pathlib import Path

import openpyxl
from openpyxl.utils import column_index_from_string as col_idx
from pptx import Presentation
from pptx.util import Inches, Pt

RAIZ = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(Path(__file__).resolve().parent))  # scripts/
sys.path.insert(0, str(RAIZ))

import avance_pptx as base  # noqa: E402  - look and feel compartido
from src.core.excel_manager import ExcelManager  # noqa: E402


# --------------------------------------------------------------------------
# Sheet layout - verified against the JULIO 2026 workbook.
#
# Hidden on purpose in `AVANCE`: DIRECTA VINOS (AE:AG), the AI/AQ spacers and
# everything from AV to BE ("Factura Presupuesto"). Hidden in `Cobertura`:
# every % and Faltan column, GABRIELA BARAKAT (W:Z) and DIRECTA VINOS (AI:AL).
# The deck shows what the sheet shows.
#
# NOTE: DIRECTA VINOS is hidden but IS included in the GFARAH column, so the
# visible vendor columns do not add up to GFARAH. That is the sheet's own
# choice, not a gap in the deck.
# --------------------------------------------------------------------------

HOJA_AVANCE = "AVANCE"
HOJA_COBERTURA = "Cobertura"

# (columna de la etiqueta en la fila 5, [Avance, %Tend, Faltan])
AVANCE_BLOQUES = [
    ("C", ["C", "D", "E"]),
    ("G", ["G", "H", "I"]),
    ("K", ["K", "L", "M"]),
    ("O", ["O", "P", "Q"]),
    ("S", ["S", "T", "U"]),
    ("W", ["W", "X", "Y"]),
    ("AA", ["AA", "AB", "AC"]),
    ("AJ", ["AJ", "AK", "AL"]),
    ("AN", ["AN", "AO", "AP"]),
    ("AS", ["AS", "AT", "AU"]),
]
AVANCE_SUBS = ["Avance", "%Tend", "Faltan"]
# `Faltan` va con un decimal porque asi lo muestra la hoja (24,6 y no 25); el
# Avance va entero, tambien como la hoja.
AVANCE_CLASES = ["num", "pct", "dec1"]

# Una sola columna visible por vendedor: el % y el Faltan estan ocultos.
COBERTURA_BLOQUES = [
    ("C", ["C"]),
    ("G", ["G"]),
    ("K", ["K"]),
    ("O", ["O"]),
    ("S", ["S"]),
    ("AA", ["AA"]),
    ("AE", ["AE"]),
    ("AM", ["AM"]),
    ("AQ", ["AQ"]),
]
COBERTURA_SUBS = ["PDV"]
COBERTURA_CLASES = ["pdv"]

FILA_VENDEDORES = 5
COL_ETIQUETA = "A"

# Secciones: (titulo, primera fila, ultima fila). El total se reconoce por su
# etiqueta, no por su numero de fila.
AVANCE_SECCIONES = [("LINEA BRANCA", 8, 22), ("OTRAS LINEAS", 23, 35)]
COBERTURA_SECCIONES = [("LINEA BRANCA", 7, 21), ("OTRAS LINEAS", 22, 37)]

MESES = ["ENERO", "FEBRERO", "MARZO", "ABRIL", "MAYO", "JUNIO", "JULIO",
         "AGOSTO", "SEPTIEMBRE", "OCTUBRE", "NOVIEMBRE", "DICIEMBRE"]


# --------------------------------------------------------------------------
# Reading the workbook
# --------------------------------------------------------------------------


def _celda(ws, fila: int, columna: str):
    return ws.cell(fila, col_idx(columna)).value


def _tiene_valores(ws) -> bool:
    """True when the sheet carries cached formula results.

    `AvancesService` writes the workbook with openpyxl, which drops every
    cached value. Without them `data_only=True` reads nothing but `None`.
    """
    columnas = [col for _etiqueta, cols in AVANCE_BLOQUES for col in cols]
    for fila in range(8, 23):
        if any(base._numero(_celda(ws, fila, c)) is not None for c in columnas):
            return True
    return False


def _abrir(archivo: Path, forzar_recalculo: bool = False):
    """Open the workbook, recalculating a COPY when the values are missing.

    The original file is never touched: LibreOffice writes into a temp dir.
    """
    libro = openpyxl.load_workbook(archivo, data_only=True)
    if not forzar_recalculo and _tiene_valores(libro[HOJA_AVANCE]):
        return libro, False

    soffice = ExcelManager._find_soffice()
    if soffice is None:
        raise RuntimeError(
            "El libro no trae valores cacheados y LibreOffice no esta instalado. "
            "Abrir el xlsx y guardarlo, o instalar libreoffice-fresh."
        )
    with tempfile.TemporaryDirectory(prefix="avance-branca-") as tmp:
        recalculado = ExcelManager._recalc_with_libreoffice(archivo, Path(tmp), soffice)
        return openpyxl.load_workbook(recalculado, data_only=True), True


def _periodo(ws) -> str:
    valor = _celda(ws, 1, "C")
    if hasattr(valor, "month"):
        return f"{MESES[valor.month - 1]} {valor.year}"
    return ""


def _etiquetas_bloques(ws, bloques: list) -> list:
    """Vendor names, read from row 5 so a rename in the sheet carries over."""
    etiquetas = []
    for columna, _cols in bloques:
        valor = _celda(ws, FILA_VENDEDORES, columna)
        etiquetas.append(str(valor).strip() if base._texto(valor) else columna)
    return etiquetas


def _filas_seccion(ws, columnas: list, desde: int, hasta: int) -> tuple:
    """(filas de detalle, fila de total) dentro de una seccion.

    Se descartan las filas sin ningun numero: la hoja tiene renglones que solo
    llevan una aclaracion en la columna A, como "(vinos fecovita)".
    """
    detalle, total = [], None
    for fila in range(desde, hasta + 1):
        etiqueta = _celda(ws, fila, COL_ETIQUETA)
        if not base._texto(etiqueta):
            continue
        etiqueta = etiqueta.strip()
        if not any(base._numero(_celda(ws, fila, c)) is not None for c in columnas):
            continue
        if etiqueta.upper().startswith("TOTAL"):
            total = (fila, etiqueta)
        else:
            detalle.append((fila, etiqueta))
    return detalle, total


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------


def _portada(prs, periodo: str, dias: tuple, origen: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base._banda(slide, 0, 0.45, base.AZUL_OSCURO)
    base._banda(slide, base.ALTO_IN - 0.45, 0.45, base.VERDE_OSCURO)

    caja = slide.shapes.add_textbox(Inches(1.1), Inches(2.5),
                                    Inches(base.ANCHO_IN - 2.2), Inches(2.6))
    marco = caja.text_frame
    marco.word_wrap = True

    parrafo = marco.paragraphs[0]
    run = parrafo.add_run()
    run.text = "AVANCE FRATELLI BRANCA"
    run.font.name = base.FUENTE
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = base.AZUL_OSCURO

    lineas = [
        (periodo, 30, True, base.AZUL_MEDIO),
        (f"Días de venta: {dias[0]}   ·   Días de avance: {dias[1]}", 15, False, base.NEGRO),
        (f"Origen: {origen.name}", 11, False, base.AZUL_MEDIO),
    ]
    for texto, size, bold, color in lineas:
        parrafo = marco.add_paragraph()
        parrafo.space_before = Pt(6)
        run = parrafo.add_run()
        run.text = str(texto)
        run.font.name = base.FUENTE
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _total_volumen(ws, filas: list, bloques: list) -> list:
    """Sum each vendor's Avance and Faltan; % = Avance / (Avance + Faltan).

    Volume adds up across categories, so this total is a sum. Coverage does
    not — see `_slide_cobertura`.
    """
    valores = []
    for _etiqueta, columnas in bloques:
        avance = base._sumar(ws, filas, columnas[0])
        faltan = base._sumar(ws, filas, columnas[2])
        cupo = None if avance is None else avance + (faltan or 0)
        valores += [avance, base._ratio(avance, cupo), faltan]
    return valores


def _slide_volumen(prs, ws, periodo: str, titulo: str, desde: int, hasta: int) -> bool:
    etiquetas = _etiquetas_bloques(ws, AVANCE_BLOQUES)
    columnas = [c for _e, cols in AVANCE_BLOQUES for c in cols]
    detalle, total = _filas_seccion(ws, columnas, desde, hasta)
    if not detalle:
        return False

    grupos = [(etiqueta, list(AVANCE_SUBS), base.AZUL_OSCURO) for etiqueta in etiquetas]
    clases = []
    for _ in AVANCE_BLOQUES:
        clases += list(AVANCE_CLASES)

    filas = [(etiqueta, base._valores(ws, fila, columnas), False) for fila, etiqueta in detalle]
    if total:
        nombre = total[1]
        filas.append((nombre, _total_volumen(ws, [f for f, _e in detalle], AVANCE_BLOQUES), True))
        subtitulo = f"{periodo}   ·   bultos y % sobre cupo"
    else:
        # Sin total a proposito: TAMBO se mide en KG y el resto en bultos.
        # Sumarlos daria un numero que no significa nada.
        subtitulo = f"{periodo}   ·   bultos y % sobre cupo   ·   sin total: TAMBO va en KG"

    slide = base._slide(prs, f"VOLUMEN — {titulo}", subtitulo)
    base._dibujar(slide, base.Tabla(grupos, "Categoría", filas, clases),
                  top_in=1.05, ancho_primera=1.5)
    return True


def _slide_cobertura(prs, ws, periodo: str, titulo: str, desde: int, hasta: int) -> bool:
    """Coverage slide. Its total row is READ from the sheet, never summed.

    Coverage counts points of sale, and a client that buys two brands is one
    PDV, not two: in JULIO 2026 the sheet's TOTAL LINEA for PABLO NAVARRO is
    201 while its brand rows add up to 352.
    """
    etiquetas = _etiquetas_bloques(ws, COBERTURA_BLOQUES)
    columnas = [c for _e, cols in COBERTURA_BLOQUES for c in cols]
    detalle, total = _filas_seccion(ws, columnas, desde, hasta)
    if not detalle:
        return False

    grupos = [(etiqueta, list(COBERTURA_SUBS), base.VERDE_OSCURO) for etiqueta in etiquetas]
    clases = []
    for _ in COBERTURA_BLOQUES:
        clases += list(COBERTURA_CLASES)

    filas = [(etiqueta, base._valores(ws, fila, columnas), False) for fila, etiqueta in detalle]
    if total:
        fila, nombre = total
        filas.append((nombre, base._valores(ws, fila, columnas), True))

    slide = base._slide(prs, f"COBERTURA — {titulo}",
                        f"{periodo}   ·   PDV cubiertos   ·   el total no es la suma de las marcas")
    base._dibujar(slide, base.Tabla(grupos, "Categoría", filas, clases),
                  top_in=1.05, ancho_primera=1.9)
    return True


# La slide de imagen es la misma en los dos decks: vive en el modulo base.
_slide_imagen = base._slide_imagen


def _capturas(archivo: Path) -> dict:
    """Los PNG que `capture_images` ya genera al lado del xlsx.

    Van al deck junto con las tablas: la tabla se lee y se filtra, la imagen es
    la hoja tal cual la ve Nahuel, con sus colores y sus guiones. Se descartan
    los `backup-*`, que son de una corrida vieja del mismo mes.
    """
    encontradas = {}
    for png in archivo.parent.glob(f"{archivo.stem}_*.png"):
        if "backup-" in png.name:
            continue
        if f"_{HOJA_AVANCE}_" in png.name:
            encontradas["AVANCE"] = png
        elif f"_{HOJA_COBERTURA}_" in png.name:
            encontradas["COBERTURA"] = png
    return encontradas


def _buscar_rechazos(periodo_iso: str) -> Path | None:
    """El PNG de rebotes del mes, ubicado por CARPETA y no por nombre.

    El nombre del archivo no sirve para elegir: `configs/rebotes.json` tiene el
    `nombre` escrito a mano ("Rebotes Junio 2026"), asi que todos los meses
    salen con el mismo nombre. La carpeta si es correcta — `service_output_dir`
    la deriva de `fecha_desde`.
    """
    carpeta = RAIZ / "data" / "output" / "reporte-rebotes" / periodo_iso
    if not carpeta.is_dir():
        return None
    pngs = sorted(carpeta.glob("*.png"), key=lambda p: p.stat().st_mtime)
    return pngs[-1] if pngs else None


# --------------------------------------------------------------------------
# Deck
# --------------------------------------------------------------------------


def _diagnostico(ws, bloques: list, secciones: list) -> list:
    """Celdas donde el TOTAL de la hoja no cierra con la suma de sus filas.

    Solo se corre sobre volumen: en cobertura el total NO es una suma y
    compararlo daria una falsa alarma en cada celda.
    """
    avisos = []
    columnas = [c for _e, cols in bloques for c in cols]
    for titulo, desde, hasta in secciones:
        detalle, total = _filas_seccion(ws, columnas, desde, hasta)
        if not total or not detalle:
            continue
        fila_total, _nombre = total
        filas = [f for f, _e in detalle]
        for _etiqueta, cols in bloques:
            for columna in (cols[0], cols[2]):  # Avance y Faltan; el % no suma
                calculado = base._sumar(ws, filas, columna)
                declarado = base._numero(_celda(ws, fila_total, columna))
                if calculado is None or declarado is None:
                    continue
                if abs(calculado - declarado) > 0.01:
                    avisos.append(
                        f"  {titulo} col {columna} fila {fila_total}: la hoja dice "
                        f"{declarado:,.2f}, la suma de sus filas da {calculado:,.2f}"
                    )
    return avisos


def poblar(prs, archivo: Path, rechazos: Path | None = None,
           diagnostico: bool = False, con_capturas: bool = True,
           portada: bool = True, con_rechazos: bool = True) -> None:
    libro, recalculado = _abrir(archivo)
    if recalculado:
        print("El libro no traia valores cacheados: se recalculo una copia con LibreOffice.")

    ws_avance = libro[HOJA_AVANCE]
    ws_cober = libro[HOJA_COBERTURA]

    periodo = _periodo(ws_avance)
    dias = (_celda(ws_avance, 1, "I"), _celda(ws_avance, 2, "I"))

    if diagnostico:
        avisos = _diagnostico(ws_avance, AVANCE_BLOQUES, AVANCE_SECCIONES)
        print("Totales de la hoja AVANCE que no cierran con sus propias filas:")
        print("\n".join(avisos) if avisos else "  (ninguno)")

    capturas = _capturas(archivo) if con_capturas else {}

    if portada:
        _portada(prs, periodo, dias, archivo)
    for titulo, desde, hasta in AVANCE_SECCIONES:
        _slide_volumen(prs, ws_avance, periodo, titulo, desde, hasta)
    if "AVANCE" in capturas:
        _slide_imagen(prs, "VOLUMEN — LA HOJA", f"{periodo}   ·   captura de `AVANCE`",
                      capturas["AVANCE"])
    for titulo, desde, hasta in COBERTURA_SECCIONES:
        _slide_cobertura(prs, ws_cober, periodo, titulo, desde, hasta)
    if "COBERTURA" in capturas:
        _slide_imagen(prs, "COBERTURA — LA HOJA", f"{periodo}   ·   captura de `Cobertura`",
                      capturas["COBERTURA"])

    # `con_rechazos=False` es el deck unificado: ahi la slide va a nivel de deck
    # porque el informe de rebotes es de toda BADIE, no de la linea Branca.
    # Poner la slide aca la dejaria dentro de la seccion equivocada, y ademas
    # repetida si el deck tambien la agrega.
    if not con_rechazos:
        return
    if rechazos and rechazos.is_file():
        _slide_imagen(prs, "RECHAZOS", f"{periodo}   ·   % de bultos rechazados por preventista",
                      rechazos)
    else:
        print("AVISO: sin imagen de rechazos, el deck sale sin esa diapositiva.")


def construir(archivo: Path, salida: Path, rechazos: Path | None = None,
              diagnostico: bool = False, con_capturas: bool = True) -> int:
    prs = base.nuevo_deck()
    poblar(prs, archivo, rechazos, diagnostico=diagnostico, con_capturas=con_capturas)
    salida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(salida)
    return len(prs.slides._sldIdLst)


def main() -> int:
    parser = argparse.ArgumentParser(description="PowerPoint del AVANCE FRATELLI BRANCA")
    parser.add_argument("--archivo", required=True,
                        help="xlsx del avance (se lee, nunca se escribe)")
    parser.add_argument("--salida", help="ruta del pptx (por defecto, junto al xlsx)")
    parser.add_argument("--rechazos", help="PNG de rebotes; por defecto se busca por periodo")
    parser.add_argument("--force", action="store_true", help="sobrescribir el pptx si ya existe")
    parser.add_argument("--diagnostico", action="store_true",
                        help="listar los totales de la hoja que no cierran")
    parser.add_argument("--sin-capturas", action="store_true",
                        help="no agregar las diapositivas con la imagen de la hoja")
    args = parser.parse_args()

    archivo = Path(args.archivo)
    if not archivo.is_file():
        print(f"ERROR: no existe {archivo}", file=sys.stderr)
        return 1

    salida = Path(args.salida) if args.salida else archivo.with_suffix(".pptx")
    if salida.exists() and not args.force:
        print(f"ERROR: {salida} ya existe. Usar --force para sobrescribir.", file=sys.stderr)
        return 1

    if args.rechazos:
        rechazos = Path(args.rechazos)
    else:
        rechazos = _buscar_rechazos(archivo.parent.name)
        if rechazos:
            print(f"Rechazos: {rechazos.name}  (carpeta {archivo.parent.name})")

    slides = construir(archivo, salida, rechazos, diagnostico=args.diagnostico,
                       con_capturas=not args.sin_capturas)
    print(f"OK: {slides} slides -> {salida}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
