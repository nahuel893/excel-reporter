#!/usr/bin/env python3
"""Build a PowerPoint deck out of an "AVANCE BADIE" workbook.

The workbook is the hand-maintained Excel that `AvancesService` refreshes in
place (see `docs/avance-badie.md`). This script only READS it: it never writes
the xlsx back.

Deck layout
-----------
1. Cover.
2. Every VOLUMEN slide, then every COBERTURA slide — never interleaved: they
   are two different readings and mixing them forces the room to jump back and
   forth. Each of the two sections opens with its own "TOTALES POR SUPERVISOR"
   slide, built from the sheet that section reads. Then, per supervisor block:
   - "VOLUMEN CERVEZAS": sheet `Avance`, three columns per category (Venta,
     % Cupo, Falta) plus the five of TOTAL CERVEZA.
   - "VOLUMEN ADO / MULTI CCU": sheet `Multicategoria`, on its own slide.
   - "COBERTURA": sheet `Cober Nueva`, opened by brand, with the same cut the
     sheet uses (Cervezas 1 and 2, Aguas, Vinos, Sidras).

Hidden columns are skipped, matching what the sheet actually shows.

Where the numbers come from
---------------------------
Two different rules, on purpose:

- **Vendor rows and the section-opening supervisor rows are READ from the
  workbook, cell by cell.** Nothing is summed, averaged or rounded — what the
  slide shows is what Nahuel sees in the Excel. Only the on-screen rendering is
  formatted.
- **The `TOTAL <codigo>` row at the foot of each detail slide is COMPUTED** as
  the sum of the vendor rows on that same slide, with percentages derived the
  way the sheet derives them (venta / cupo, PDV / OBJ). It has to be: it totals
  exactly the rows above it, so reading a workbook cell there could contradict
  the very rows it sits under.

That distinction matters because several of the workbook's own total rows are
stale, each in a different place (checked against JULIO 2026):

- `Avance`, summary band, FGUANTAY: TOTAL CERVEZA 22.594,57 instead of
  23.340,65 — its SUM range never grew when LORENA TARITOLAY was added.
- `Cober Nueva`, VCHAPUR block total: SIDRAS Y LICORES 23 PDV instead of 21 —
  the row carries FGUANTAY's figure.

The deck does not silently pick a winner. `--diagnostico` prints every cell
where the workbook's own total row disagrees with the sum of its vendors, so the
mismatch gets fixed in the Excel instead of being papered over here.

Usage
-----
    python scripts/avance_pptx.py \
        --archivo "data/output/avances/2026-07/AVANCE BADIE - JULIO 2026.xlsx"
"""
from __future__ import annotations

import argparse
import sys
from dataclasses import dataclass, field
from pathlib import Path

import openpyxl
from openpyxl.utils import column_index_from_string as col_idx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# --------------------------------------------------------------------------
# Sheet layout - verified against the JULIO 2026 workbook.
#
# `Avance` hides SALTA CAUTIVA1 (G:I) and everything from Y to AL, so those
# categories are left out on purpose: the deck shows what the sheet shows.
# `Multicategoria` hides K:V (ALIMENTOS & JUGOS, FERNET, VINOS, LA HUERTA),
# which are also broken with #REF! in the source workbook.
# --------------------------------------------------------------------------

# (etiqueta, venta, %, falta) — las tres columnas que la hoja da por categoria.
# El cupo no es columna propia: la hoja escribe Falta = Cupo - Venta, asi que
# Cupo = Venta + Falta.
AVANCE_CATEGORIAS = [
    ("SALTA", "D", "E", "F"),
    ("HEINEKEN", "J", "K", "L"),
    ("IMPERIAL", "M", "N", "O"),
    ("MILLER", "P", "Q", "R"),
    ("MULTICERV.", "S", "T", "U"),
    ("IMPORTADAS", "V", "W", "X"),
]
CATEGORIA_SUBS = ["Venta", "% Cupo", "Falta"]
CATEGORIA_CLASES = ["num", "pct", "num"]

# (etiqueta, columna, clase). AR ("Vta. Diaria p/ Cupo") queda afuera: en el
# libro es #DIV/0! en todas las filas.
AVANCE_TOTAL_ETIQUETA = "TOTAL CERVEZA"
AVANCE_TOTAL_COLS = [
    ("Venta", "AM", "num"),
    ("Cupo", "AN", "num"),
    ("Tend.", "AO", "num"),
    ("% Tend.", "AP", "pct"),
    ("Dif.", "AQ", "num"),
]
AVANCE_TOTAL_VENTA = "AM"
AVANCE_TOTAL_CUPO = "AN"

MULTI_CATEGORIAS = [
    ("AGUAS DANONE", "D", "E", "F"),
    ("MULTI CCU", "H", "I", "J"),
]

# Totales por generico, para el resumen. (etiqueta, PDV, OBJ, Faltan, %)
COBER_GENERICOS = [
    ("CERVEZAS", "AT", "AU", "AV", "AW"),
    ("AGUAS DANONE", "BU", "BV", "BW", "BX"),
    ("VINOS CCU", "CT", "CU", "CV", "CW"),
    ("SIDRAS Y LICORES", "DS", "DT", "DU", "DV"),
]

# Detalle por marca. Un slide por bloque, con el mismo corte que usa la hoja:
# CERVEZAS son 10 marcas y no entran legibles en una sola diapositiva, asi que
# van separadas igual que en `Cober Nueva` (Cervezas 1 y Cervezas 2).
# Las marcas de VINOS y SIDRAS traen tres columnas, no cuatro: la hoja no les
# calcula "Faltan". El TOTAL de cada generico si las trae.
COBER_SLIDES = [
    ("COBERTURA CERVEZAS 1", [
        ("SALTA", ["C", "D", "E", "F"]),
        ("HEINEKEN", ["G", "H", "I", "J"]),
        ("IMPERIAL", ["K", "L", "M", "N"]),
        ("MILLER", ["O", "P", "Q", "R"]),
    ]),
    ("COBERTURA CERVEZAS 2", [
        ("BIECKERT", ["V", "W", "X", "Y"]),
        ("SCHNEIDER", ["Z", "AA", "AB", "AC"]),
        ("AMSTEL", ["AD", "AE", "AF", "AG"]),
        ("KUNSTMAN", ["AH", "AI", "AJ", "AK"]),
        ("BLUE MOON", ["AL", "AM", "AN", "AO"]),
        ("SALTA CAUTIVA1", ["AP", "AQ", "AR", "AS"]),
        ("TOTAL CERVEZAS", ["AT", "AU", "AV", "AW"]),
    ]),
    ("COBERTURA AGUAS DANONE", [
        ("LEVITE", ["BA", "BB", "BC", "BD"]),
        ("VILLAVICENCIO", ["BE", "BF", "BG", "BH"]),
        ("VILLA DEL SUR", ["BI", "BJ", "BK", "BL"]),
        ("BRIO", ["BM", "BN", "BO", "BP"]),
        ("FULL SPORT", ["BQ", "BR", "BS", "BT"]),
        ("TOTAL AGUAS", ["BU", "BV", "BW", "BX"]),
    ]),
    ("COBERTURA VINOS CCU", [
        ("COLON", ["CB", "CC", "CD"]),
        ("LA CELIA", ["CE", "CF", "CG"]),
        ("GRAFFIGNA", ["CH", "CI", "CJ"]),
        ("EUGENIO BUSTOS", ["CK", "CL", "CM"]),
        ("O-61", ["CN", "CO", "CP"]),
        ("SANTA SILVIA", ["CQ", "CR", "CS"]),
        ("TOTAL VINOS", ["CT", "CU", "CV", "CW"]),
    ]),
    ("COBERTURA SIDRAS Y LICORES", [
        ("REAL", ["DA", "DB", "DC"]),
        ("LA VICTORIA", ["DD", "DE", "DF"]),
        ("SAENZ BRIONES", ["DG", "DH", "DI"]),
        ("EL ABUELO", ["DJ", "DK", "DL"]),
        ("PEHUENIA", ["DM", "DN", "DO"]),
        ("MISTRAL", ["DP", "DQ", "DR"]),
        ("TOTAL SIDRAS", ["DS", "DT", "DU", "DV"]),
    ]),
]

COBER_SUBS = {3: ["PDV", "OBJ", "%"], 4: ["PDV", "OBJ", "Faltan", "%"]}
COBER_CLASES = {3: ["pdv", "obj", "pct"], 4: ["pdv", "obj", "obj", "pct"]}

ENCABEZADOS_IGNORADOS = {"vendedor", "super", "cobertura", "total", "n"}

# --------------------------------------------------------------------------
# Look and feel
# --------------------------------------------------------------------------

FUENTE = "Calibri"
ANCHO_IN = 13.333
ALTO_IN = 7.5

AZUL_OSCURO = RGBColor(0x1F, 0x38, 0x64)
AZUL_MEDIO = RGBColor(0x2E, 0x5C, 0x9A)
# Cobertura keeps its own colour so a mixed slide reads at a glance.
VERDE_OSCURO = RGBColor(0x1B, 0x4D, 0x4A)
VERDE_MEDIO = RGBColor(0x2C, 0x7A, 0x74)
GRIS_FILA = RGBColor(0xF2, 0xF4, 0xF8)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO = RGBColor(0x25, 0x25, 0x25)
GRIS_TOTAL = RGBColor(0xD9, 0xDF, 0xEA)

VERDE = RGBColor(0x1E, 0x7B, 0x34)
AMBAR = RGBColor(0xB8, 0x6E, 0x00)
ROJO = RGBColor(0xB3, 0x1B, 0x1B)

UMBRAL_OK = 1.0
UMBRAL_ALERTA = 0.9


# --------------------------------------------------------------------------
# Extraction
# --------------------------------------------------------------------------


def _texto(valor) -> bool:
    return isinstance(valor, str) and valor.strip() != ""


def _numero(valor):
    """Return the value only when it is a usable number.

    Cached formula errors (#REF!, #DIV/0!) come back as strings and must not be
    mistaken for data.
    """
    if isinstance(valor, bool):
        return None
    return valor if isinstance(valor, (int, float)) else None


@dataclass
class Bloque:
    """A supervisor band: its vendors, in sheet order."""

    codigo: str
    vendedores: list = field(default_factory=list)
    filas: dict = field(default_factory=dict)  # vendedor -> row


def _miembros_por_codigo(ws, col_nombre: str, col_codigo: str, max_row: int) -> dict:
    """code -> [(nombre, row), ...] for every row that looks like a member."""
    ci_nombre, ci_codigo = col_idx(col_nombre), col_idx(col_codigo)
    miembros = {}
    for row in range(1, max_row + 1):
        nombre = ws.cell(row, ci_nombre).value
        codigo = ws.cell(row, ci_codigo).value
        if not _texto(nombre) or not _texto(codigo):
            continue
        nombre, codigo = nombre.strip(), codigo.strip()
        if nombre.lower() in ENCABEZADOS_IGNORADOS or nombre.lower().startswith("cobertura"):
            continue
        if codigo.lower() in ENCABEZADOS_IGNORADOS:
            continue
        miembros.setdefault(codigo, []).append((nombre, row))
    return miembros


def _bloques(ws, col_nombre: str, col_codigo: str, max_row: int) -> list:
    """Group vendor rows by supervisor code.

    A supervisor block is a code whose members are all plain vendor names. The
    roll-up code (GFARAH) is excluded because its members are codes themselves,
    and single-line bands (SUB DISTRIBUIDOR, DIRECTA) only reach the summary.
    """
    miembros = _miembros_por_codigo(ws, col_nombre, col_codigo, max_row)
    codigos = set(miembros)
    bloques = []
    for codigo, filas in miembros.items():
        nombres = [n for n, _ in filas]
        if any(n in codigos for n in nombres) or len(nombres) < 2:
            continue
        bloque = Bloque(codigo=codigo)
        for nombre, row in filas:
            bloque.vendedores.append(nombre)
            bloque.filas[nombre] = row
        bloques.append(bloque)
    return bloques


def _codigo_rollup(ws, col_nombre: str, col_codigo: str, max_row: int):
    """The code whose members are supervisor codes (GFARAH in this workbook)."""
    miembros = _miembros_por_codigo(ws, col_nombre, col_codigo, max_row)
    codigos = set(miembros)
    for codigo, filas in miembros.items():
        if any(n in codigos for n, _ in filas):
            return codigo
    return None


def _filas_resumen(ws, col_codigo: str, col_nombre: str, max_row: int,
                   codigos: set, exigir_nombre_vacio: bool) -> dict:
    """Map supervisor code -> row inside the sheet's own summary band.

    Used only by `--diagnostico`: the deck never takes a total from the
    workbook, it sums the rows it shows.
    """
    ci_codigo, ci_nombre = col_idx(col_codigo), col_idx(col_nombre)
    encontrado = {}
    for row in range(1, max_row + 1):
        codigo = ws.cell(row, ci_codigo).value
        if not _texto(codigo) or codigo.strip() not in codigos:
            continue
        if exigir_nombre_vacio and _texto(ws.cell(row, ci_nombre).value):
            continue
        encontrado[codigo.strip()] = row
    return encontrado


def _filas_por_supervisor(ws, col_nombre: str, col_codigo: str, max_row: int,
                          columnas: list, col_banda: str = "B",
                          col_vacia: str = "A",
                          exigir_nombre_vacio: bool = False) -> list:
    """Una fila por supervisor, LEIDA de la banda de resumen del libro.

    Estas son las filas que abren cada seccion. El valor sale de la celda tal
    cual: no se suman los vendedores, no se recalcula ningun porcentaje y no se
    redondea nada. Lo que se ve en la slide es lo que Nahuel ve en el Excel.

    La banda vive al pie de la hoja y NO tiene la misma forma en las tres:

    - `Avance` (52-55): el codigo va en B, la columna A queda vacia.
    - `Cober Nueva` (51-54): el codigo tambien va en B, pero A trae un 0
      numerico, no un vacio.

    Por eso la busqueda va por `_filas_resumen`, que es el mismo helper que usa
    `--diagnostico` y ya contempla esa diferencia via `exigir_nombre_vacio`.
    Buscar "filas cuyo nombre es un codigo" no sirve: en `Cober Nueva` el 0 de
    la columna A no es texto y la banda entera se pierde.

    Se filtran contra los bloques reales para dejar afuera las bandas de una
    sola linea (DIRECTA, SUB DISTRIBUIDOR), que no son supervisores.

    Cuando la hoja y la suma de sus propias filas no cierran, eso lo reporta
    `--diagnostico`; el deck no elige un ganador por su cuenta.
    """
    supervisores = {b.codigo for b in _bloques(ws, col_nombre, col_codigo, max_row)}
    if not supervisores:
        return []
    banda = _filas_resumen(ws, col_banda, col_vacia, max_row, supervisores,
                           exigir_nombre_vacio)
    return [(codigo, _valores(ws, row, columnas), False)
            for codigo, row in sorted(banda.items(), key=lambda par: par[1])]


def _valores(ws, row, columnas: list) -> list:
    if not row:
        return [None] * len(columnas)
    return [_numero(ws.cell(row, col_idx(c)).value) for c in columnas]


def _sumar(ws, filas: list, columna: str):
    """Sum one column over several rows. None when no row holds a number."""
    ci = col_idx(columna)
    total, visto = 0.0, False
    for row in filas:
        valor = _numero(ws.cell(row, ci).value)
        if valor is not None:
            total += valor
            visto = True
    return total if visto else None


def _ratio(numerador, denominador):
    if numerador is None or not denominador:
        return None
    return numerador / denominador


# --------------------------------------------------------------------------
# Column layouts. Each builder returns the value list for one table row, in the
# same order as the headers declared by the matching slide.
# --------------------------------------------------------------------------


def _cols_volumen_avance() -> list:
    cols = [c for _, venta, pct, falta in AVANCE_CATEGORIAS for c in (venta, pct, falta)]
    return cols + [col for _etiqueta, col, _clase in AVANCE_TOTAL_COLS]


def _cols_volumen_multi() -> list:
    return [c for _, venta, pct, falta in MULTI_CATEGORIAS for c in (venta, pct, falta)]


def _cols_cobertura_generico() -> list:
    """Totales por generico — los que van al resumen, sin abrir por marca."""
    return [c for _, *cols in COBER_GENERICOS for c in cols]


# Only these add up across rows: a percentage column is a ratio, not a sum.
def _cols_aditivas_avance() -> list:
    cols = [c for _, venta, _pct, falta in AVANCE_CATEGORIAS for c in (venta, falta)]
    return cols + [col for _etiqueta, col, clase in AVANCE_TOTAL_COLS if clase != "pct"]


def _cols_aditivas_multi() -> list:
    return [c for _, venta, _pct, falta in MULTI_CATEGORIAS for c in (venta, falta)]


def _cols_aditivas_cobertura() -> list:
    """Todas las de los slides de marca menos el %: PDV, OBJ y Faltan suman."""
    cols = []
    for _titulo, bloques in COBER_SLIDES:
        for _marca, columnas in bloques:
            cols += columnas[:-1]  # la ultima siempre es el %
    return cols


def _total_categorias(ws, filas: list, categorias: list) -> list:
    """[venta, %, falta] per category, summed. % = venta / (venta + falta)."""
    valores = []
    for _etiqueta, venta_col, _pct_col, falta_col in categorias:
        venta = _sumar(ws, filas, venta_col)
        falta = _sumar(ws, filas, falta_col)
        cupo = None if venta is None else venta + (falta or 0)
        valores += [venta, _ratio(venta, cupo), falta]
    return valores


def _total_volumen_avance(ws, filas: list) -> list:
    valores = _total_categorias(ws, filas, AVANCE_CATEGORIAS)
    venta = _sumar(ws, filas, AVANCE_TOTAL_VENTA)
    cupo = _sumar(ws, filas, AVANCE_TOTAL_CUPO)
    for _etiqueta, col, clase in AVANCE_TOTAL_COLS:
        valores.append(_ratio(venta, cupo) if clase == "pct" else _sumar(ws, filas, col))
    return valores


def _total_volumen_multi(ws, filas: list) -> list:
    return _total_categorias(ws, filas, MULTI_CATEGORIAS)


def _total_cobertura(ws, filas: list) -> list:
    """[PDV, OBJ, Faltan, %] per generic, summed over `filas`."""
    valores = []
    for _etiqueta, pdv_col, obj_col, falta_col, _pct_col in COBER_GENERICOS:
        pdv = _sumar(ws, filas, pdv_col)
        obj = _sumar(ws, filas, obj_col)
        valores += [pdv, obj, _sumar(ws, filas, falta_col), _ratio(pdv, obj)]
    return valores


# --------------------------------------------------------------------------
# Formatting
# --------------------------------------------------------------------------


def _fmt_numero(valor, decimales: int) -> str:
    entero, _, decimal = f"{valor:,.{decimales}f}".partition(".")
    entero = entero.replace(",", ".")
    return f"{entero},{decimal}" if decimal else entero


def _fmt(valor, clase: str) -> str:
    """Display only. The stored value is never rounded, just rendered.

    - `num` (venta, cupo, falta): sin decimales, son bultos.
    - `pct`: con un decimal, porque la diferencia entre 96,4% y 96,9% importa.
    - `pdv`: PDV son enteros.
    - `obj` / `dec1`: un decimal. El objetivo de cobertura es fraccionario y en
      SIDRAS anda por 9,3, asi que redondearlo lo dejaria sin significado; lo
      mismo vale para cualquier columna que la hoja muestre con un decimal.
    """
    if valor is None:
        return "-"
    if clase == "pct":
        # Con objetivos casi en cero el % se dispara (AMSTEL, 1.027%). Ahi el
        # decimal no aporta y ademas parte la celda en dos lineas.
        puntos = valor * 100
        return f"{_fmt_numero(puntos, 0 if abs(puntos) >= 1000 else 1)}%"
    if clase in ("obj", "dec1"):
        return _fmt_numero(valor, 1)
    return _fmt_numero(valor, 0)


def _color_pct(valor):
    if valor is None:
        return None
    if valor >= UMBRAL_OK:
        return VERDE
    if valor >= UMBRAL_ALERTA:
        return AMBAR
    return ROJO


# --------------------------------------------------------------------------
# Slide building
# --------------------------------------------------------------------------


def _run(celda, texto: str, size: int, bold: bool, color: RGBColor,
         align=PP_ALIGN.RIGHT, wrap: bool = True, margen: float = 0.03) -> None:
    """`wrap=False` y margen chico para las celdas de dato.

    Un numero partido en dos lineas ("142,5" arriba y "%" abajo) se lee mucho
    peor que uno que roza el borde, y ademas estira toda la fila. Los
    encabezados si envuelven: ahi el salto de linea ayuda.
    """
    celda.text_frame.clear()
    celda.text_frame.word_wrap = wrap
    parrafo = celda.text_frame.paragraphs[0]
    parrafo.alignment = align
    run = parrafo.add_run()
    run.text = texto
    run.font.name = FUENTE
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    celda.vertical_anchor = MSO_ANCHOR.MIDDLE
    celda.margin_left = Inches(margen)
    celda.margin_right = Inches(margen)
    celda.margin_top = 0
    celda.margin_bottom = 0


def _pintar(celda, color: RGBColor) -> None:
    celda.fill.solid()
    celda.fill.fore_color.rgb = color


def _slide(prs, titulo: str, subtitulo: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    caja = slide.shapes.add_textbox(Inches(0.35), Inches(0.22), Inches(ANCHO_IN - 0.7), Inches(0.55))
    parrafo = caja.text_frame.paragraphs[0]
    run = parrafo.add_run()
    run.text = titulo
    run.font.name = FUENTE
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = AZUL_OSCURO
    if subtitulo:
        sub = caja.text_frame.add_paragraph()
        run = sub.add_run()
        run.text = subtitulo
        run.font.name = FUENTE
        run.font.size = Pt(12)
        run.font.color.rgb = AZUL_MEDIO
    return slide


@dataclass
class Tabla:
    """A two-header-row table: groups on top, sub-columns below."""

    grupos: list          # [(group label, [sub labels], dark colour)]
    etiqueta_primera: str
    filas: list           # [(nombre, valores, es_total)]
    clases: list          # one per value column


def _dibujar(slide, tabla: Tabla, top_in: float, ancho_primera: float) -> None:
    subcols = [s for _, subs, _color in tabla.grupos for s in subs]
    n_cols = 1 + len(subcols)
    n_filas = 2 + len(tabla.filas)

    alto_disponible = ALTO_IN - top_in - 0.35
    alto_fila = min(0.45, alto_disponible / n_filas)
    alto_tabla = alto_fila * n_filas
    top_in += max(0.0, (alto_disponible - alto_tabla) / 2)

    forma = slide.shapes.add_table(
        n_filas, n_cols,
        Inches(0.35), Inches(top_in),
        Inches(ANCHO_IN - 0.7), Inches(alto_tabla),
    )
    tabla_pptx = forma.table
    tabla_pptx.first_row = False
    tabla_pptx.horz_banding = False

    ancho_resto = (ANCHO_IN - 0.7 - ancho_primera) / len(subcols)
    tabla_pptx.columns[0].width = Inches(ancho_primera)
    for i in range(1, n_cols):
        tabla_pptx.columns[i].width = Inches(ancho_resto)
    for fila in tabla_pptx.rows:
        fila.height = Inches(alto_fila)

    # El escalon de 28 sale de medir: con 30 subcolumnas un "142,5%" a 7 pt no
    # entra en 0,37" y LibreOffice lo parte en dos lineas.
    if len(subcols) <= 14:
        fuente_dato, fuente_cab = 9, 9
    elif len(subcols) <= 24:
        fuente_dato, fuente_cab = 8, 7
    elif len(subcols) <= 28:
        fuente_dato, fuente_cab = 7, 6
    else:
        fuente_dato, fuente_cab = 6.5, 6

    celda = tabla_pptx.cell(0, 0)
    celda.merge(tabla_pptx.cell(1, 0))
    _pintar(celda, tabla.grupos[0][2])
    _run(celda, tabla.etiqueta_primera, fuente_cab + 1, True, BLANCO, PP_ALIGN.LEFT)

    col = 1
    for grupo, subs, color in tabla.grupos:
        claro = VERDE_MEDIO if color == VERDE_OSCURO else AZUL_MEDIO
        primera, ultima = col, col + len(subs) - 1
        origen = tabla_pptx.cell(0, primera)
        if ultima > primera:
            origen.merge(tabla_pptx.cell(0, ultima))
        _pintar(origen, color)
        _run(origen, grupo, fuente_cab, True, BLANCO, PP_ALIGN.CENTER)
        for offset, sub in enumerate(subs):
            hijo = tabla_pptx.cell(1, primera + offset)
            _pintar(hijo, claro)
            _run(hijo, sub, fuente_cab, False, BLANCO, PP_ALIGN.CENTER)
        col = ultima + 1

    for i, (nombre, valores, es_total) in enumerate(tabla.filas):
        fila = 2 + i
        fondo = GRIS_TOTAL if es_total else (GRIS_FILA if i % 2 else BLANCO)
        celda = tabla_pptx.cell(fila, 0)
        _pintar(celda, fondo)
        _run(celda, nombre, fuente_dato, es_total, NEGRO, PP_ALIGN.LEFT)
        for j, valor in enumerate(valores):
            clase = tabla.clases[j]
            celda = tabla_pptx.cell(fila, 1 + j)
            _pintar(celda, fondo)
            color = _color_pct(valor) if clase == "pct" else None
            _run(celda, _fmt(valor, clase), fuente_dato, es_total, color or NEGRO,
                 wrap=False, margen=0.015)


# --------------------------------------------------------------------------
# Deck
# --------------------------------------------------------------------------

MESES = ["ENERO", "FEBRERO", "MARZO", "ABRIL", "MAYO", "JUNIO", "JULIO",
         "AGOSTO", "SEPTIEMBRE", "OCTUBRE", "NOVIEMBRE", "DICIEMBRE"]


def _periodo(ws_avance) -> str:
    valor = ws_avance.cell(1, col_idx("B")).value
    if hasattr(valor, "month"):
        return f"{MESES[valor.month - 1]} {valor.year}"
    return ""


def _banda(slide, top_in: float, alto_in: float, color: RGBColor) -> None:
    forma = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(top_in), Inches(ANCHO_IN), Inches(alto_in))
    forma.fill.solid()
    forma.fill.fore_color.rgb = color
    forma.line.fill.background()
    forma.shadow.inherit = False


def _portada(prs, periodo: str, dias: tuple, origen: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _banda(slide, 0, 0.45, AZUL_OSCURO)
    _banda(slide, ALTO_IN - 0.45, 0.45, VERDE_OSCURO)

    caja = slide.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(ANCHO_IN - 2.2), Inches(2.6))
    marco = caja.text_frame
    marco.word_wrap = True

    parrafo = marco.paragraphs[0]
    run = parrafo.add_run()
    run.text = "AVANCE BADIE"
    run.font.name = FUENTE
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = AZUL_OSCURO

    lineas = [
        (periodo, 30, True, AZUL_MEDIO),
        (f"Días hábiles: {dias[0]}   ·   Días de venta: {dias[1]}   ·   Faltan: {dias[2]}",
         15, False, NEGRO),
        (f"Origen: {origen.name}", 11, False, AZUL_MEDIO),
    ]
    for texto, size, bold, color in lineas:
        parrafo = marco.add_paragraph()
        parrafo.space_before = Pt(6)
        run = parrafo.add_run()
        run.text = str(texto)
        run.font.name = FUENTE
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _cabecera_volumen_cervezas() -> tuple:
    """(grupos, clases) de la tabla de cervezas.

    La comparte el detalle por vendedor con la apertura por supervisor: si
    manana entra una categoria nueva, las dos tablas la toman juntas.
    """
    grupos = [(nombre, list(CATEGORIA_SUBS), AZUL_OSCURO) for nombre, *_ in AVANCE_CATEGORIAS]
    grupos.append((AVANCE_TOTAL_ETIQUETA,
                   [etiqueta for etiqueta, _col, _clase in AVANCE_TOTAL_COLS],
                   AZUL_OSCURO))

    clases = []
    for _ in AVANCE_CATEGORIAS:
        clases += list(CATEGORIA_CLASES)
    clases += [clase for _etiqueta, _col, clase in AVANCE_TOTAL_COLS]
    return grupos, clases


def _cabecera_volumen_multi() -> tuple:
    grupos = [(nombre, list(CATEGORIA_SUBS), AZUL_OSCURO) for nombre, *_ in MULTI_CATEGORIAS]
    clases = []
    for _ in MULTI_CATEGORIAS:
        clases += list(CATEGORIA_CLASES)
    return grupos, clases


def _cabecera_cobertura(bloques: list) -> tuple:
    grupos = [(marca, COBER_SUBS[len(columnas)], VERDE_OSCURO) for marca, columnas in bloques]
    clases = []
    for _marca, columnas in bloques:
        clases += COBER_CLASES[len(columnas)]
    return grupos, clases


def _slide_volumen_cervezas(prs, codigo: str, periodo: str, ws_avance,
                            bloque: Bloque) -> None:
    grupos, clases = _cabecera_volumen_cervezas()

    columnas = _cols_volumen_avance()
    filas = [(v, _valores(ws_avance, bloque.filas[v], columnas), False)
             for v in bloque.vendedores]
    filas_bloque = [bloque.filas[v] for v in bloque.vendedores]
    filas.append((f"TOTAL {codigo}", _total_volumen_avance(ws_avance, filas_bloque), True))

    slide = _slide(prs, f"VOLUMEN CERVEZAS — {codigo}",
                   f"{periodo}   ·   bultos y % sobre cupo")
    _dibujar(slide, Tabla(grupos, "Vendedor", filas, clases), top_in=1.05, ancho_primera=1.6)


def _slide_volumen_multi(prs, codigo: str, periodo: str, ws_multi,
                         bloque: Bloque) -> None:
    grupos = [(nombre, list(CATEGORIA_SUBS), AZUL_OSCURO) for nombre, *_ in MULTI_CATEGORIAS]
    clases = []
    for _ in MULTI_CATEGORIAS:
        clases += list(CATEGORIA_CLASES)

    columnas = _cols_volumen_multi()
    filas = [(v, _valores(ws_multi, bloque.filas[v], columnas), False)
             for v in bloque.vendedores]
    filas_bloque = [bloque.filas[v] for v in bloque.vendedores]
    filas.append((f"TOTAL {codigo}", _total_volumen_multi(ws_multi, filas_bloque), True))

    slide = _slide(prs, f"VOLUMEN ADO / MULTI CCU — {codigo}",
                   f"{periodo}   ·   bultos y % sobre cupo")
    _dibujar(slide, Tabla(grupos, "Vendedor", filas, clases), top_in=1.05, ancho_primera=2.4)


def _total_marcas(ws, filas: list, bloques: list) -> list:
    """Suma PDV/OBJ/Faltan de cada marca; el % se deriva como PDV / OBJ."""
    valores = []
    for _marca, columnas in bloques:
        sumas = [_sumar(ws, filas, c) for c in columnas[:-1]]
        valores += sumas + [_ratio(sumas[0], sumas[1])]
    return valores


def _slide_cobertura(prs, titulo: str, bloques: list, codigo: str, periodo: str,
                     ws_cober, bloque: Bloque) -> None:
    grupos = [(marca, COBER_SUBS[len(columnas)], VERDE_OSCURO) for marca, columnas in bloques]
    clases = []
    for _marca, columnas in bloques:
        clases += COBER_CLASES[len(columnas)]
    columnas = [c for _marca, cols in bloques for c in cols]

    filas = [(v, _valores(ws_cober, bloque.filas[v], columnas), False) for v in bloque.vendedores]
    filas_bloque = [bloque.filas[v] for v in bloque.vendedores]
    filas.append((f"TOTAL {codigo}", _total_marcas(ws_cober, filas_bloque, bloques), True))

    slide = _slide(prs, f"{titulo} — {codigo}",
                   f"{periodo}   ·   PDV cubiertos sobre objetivo")
    _dibujar(slide, Tabla(grupos, "Vendedor", filas, clases), top_in=1.05, ancho_primera=1.55)


def _slide_totales(prs, titulo: str, subtitulo: str, grupos: list, clases: list,
                   filas: list, ancho_primera: float) -> None:
    """Apertura de seccion: los supervisores y sus totales, nada mas."""
    if not filas:
        return
    slide = _slide(prs, titulo, subtitulo)
    _dibujar(slide, Tabla(grupos, "Supervisor", filas, clases),
             top_in=1.05, ancho_primera=ancho_primera)


def _slide_totales_volumen_cervezas(prs, periodo: str, ws_avance) -> None:
    grupos, clases = _cabecera_volumen_cervezas()
    filas = _filas_por_supervisor(ws_avance, "B", "C", ws_avance.max_row,
                                  _cols_volumen_avance())
    _slide_totales(prs, "VOLUMEN CERVEZAS — TOTALES POR SUPERVISOR",
                   f"{periodo}   ·   bultos y % sobre cupo",
                   grupos, clases, filas, ancho_primera=1.6)


def _slide_totales_volumen_multi(prs, periodo: str, ws_multi) -> None:
    grupos, clases = _cabecera_volumen_multi()
    filas = _filas_por_supervisor(ws_multi, "B", "C", ws_multi.max_row,
                                  _cols_volumen_multi())
    _slide_totales(prs, "VOLUMEN ADO / MULTI CCU — TOTALES POR SUPERVISOR",
                   f"{periodo}   ·   bultos y % sobre cupo",
                   grupos, clases, filas, ancho_primera=2.4)


def _slide_totales_cobertura(prs, titulo: str, bloques: list, periodo: str,
                             ws_cober) -> None:
    grupos, clases = _cabecera_cobertura(bloques)
    columnas = [c for _marca, cols in bloques for c in cols]
    # En `Cober Nueva` la columna A de la banda trae un 0, no un vacio, y las
    # filas de vendedor sí llevan el codigo en B: sin este flag la banda se
    # confunde con la ultima fila de vendedor de cada supervisor.
    filas = _filas_por_supervisor(ws_cober, "A", "B", ws_cober.max_row, columnas,
                                  exigir_nombre_vacio=True)
    _slide_totales(prs, f"{titulo} — TOTALES POR SUPERVISOR",
                   f"{periodo}   ·   PDV cubiertos sobre objetivo",
                   grupos, clases, filas, ancho_primera=1.55)


def _otros(ws, col_nombre: str, col_codigo: str, max_row: int, codigos_bloque: set) -> list:
    """Single-line bands: DIRECTA and SUB DISTRIBUIDOR, in sheet order.

    They are part of the grand total but have no supervisor slide of their own,
    so they only show up on the summary.
    """
    miembros = _miembros_por_codigo(ws, col_nombre, col_codigo, max_row)
    codigos = set(miembros)
    sueltos = []
    for codigo, filas in miembros.items():
        if codigo in codigos_bloque:
            continue
        for nombre, row in filas:
            if nombre not in codigos:
                sueltos.append((nombre, row))
    return sorted(sueltos, key=lambda par: par[1])


def _diagnostico(ws, etiqueta: str, bloque: Bloque, columnas: list,
                 fila_declarada, contexto: str) -> list:
    """Cells where a total row of the sheet disagrees with the sum of its rows."""
    if not fila_declarada:
        return []
    filas = [bloque.filas[v] for v in bloque.vendedores]
    avisos = []
    for columna in columnas:
        calculado = _sumar(ws, filas, columna)
        declarado = _numero(ws.cell(fila_declarada, col_idx(columna)).value)
        if calculado is None or declarado is None:
            continue
        if abs(calculado - declarado) > 0.01:
            avisos.append(
                f"  {etiqueta} / {contexto} / {bloque.codigo}: col {columna} fila "
                f"{fila_declarada} dice {declarado:,.2f}, la suma de sus filas da "
                f"{calculado:,.2f}"
            )
    return avisos


def _avisos_totales(ws, etiqueta: str, bloques: list, columnas: list,
                    resumen: dict) -> list:
    avisos = []
    for bloque in bloques:
        fila_bloque = max(bloque.filas.values()) + 1
        avisos += _diagnostico(ws, etiqueta, bloque, columnas, fila_bloque, "total del bloque")
        avisos += _diagnostico(ws, etiqueta, bloque, columnas,
                               resumen.get(bloque.codigo), "banda de resumen")
    return avisos


def construir(archivo: Path, salida: Path, diagnostico: bool = False) -> int:
    libro = openpyxl.load_workbook(archivo, data_only=True)
    ws_avance = libro["Avance"]
    ws_multi = libro["Multicategoria"]
    ws_cober = libro["Cober Nueva"]

    periodo = _periodo(ws_avance)
    dias = tuple(ws_avance.cell(r, col_idx("AN")).value for r in (1, 2, 3))

    bloques_av = _bloques(ws_avance, "B", "C", ws_avance.max_row)
    bloques_mu = {b.codigo: b for b in _bloques(ws_multi, "B", "C", ws_multi.max_row)}
    bloques_cb = {b.codigo: b for b in _bloques(ws_cober, "A", "B", ws_cober.max_row)}
    codigos_bloque = {b.codigo for b in bloques_av}

    if diagnostico:
        rollup = _codigo_rollup(ws_avance, "B", "C", ws_avance.max_row)
        codigos = set(_miembros_por_codigo(ws_avance, "B", "C", ws_avance.max_row))
        if rollup:
            codigos.add(rollup)
        avisos = _avisos_totales(
            ws_avance, "Avance", bloques_av, _cols_aditivas_avance(),
            _filas_resumen(ws_avance, "B", "A", ws_avance.max_row, codigos, False))
        avisos += _avisos_totales(
            ws_multi, "Multicategoria", list(bloques_mu.values()), _cols_aditivas_multi(),
            _filas_resumen(ws_multi, "B", "A", ws_multi.max_row, codigos, False))
        avisos += _avisos_totales(
            ws_cober, "Cober Nueva", list(bloques_cb.values()), _cols_aditivas_cobertura(),
            _filas_resumen(ws_cober, "B", "A", ws_cober.max_row, codigos, True))
        print("Filas de total del libro que no cierran con sus propias filas:")
        print("\n".join(avisos) if avisos else "  (ninguna)")


    prs = Presentation()
    prs.slide_width = Inches(ANCHO_IN)
    prs.slide_height = Inches(ALTO_IN)

    _portada(prs, periodo, dias, archivo)

    # Primero todo el volumen, despues toda la cobertura: son dos lecturas
    # distintas y mezclarlas obliga a saltar de una a otra en la reunion.
    #
    # Cada seccion abre con los totales por supervisor de SU hoja, leidos de la
    # banda de resumen del libro. Antes habia una sola slide de resumen que
    # ponia volumen y cobertura en la misma tabla: eso era justamente el salto
    # que se queria evitar.
    _slide_totales_volumen_cervezas(prs, periodo, ws_avance)
    _slide_totales_volumen_multi(prs, periodo, ws_multi)
    for bloque in bloques_av:
        _slide_volumen_cervezas(prs, bloque.codigo, periodo, ws_avance, bloque)
        bloque_mu = bloques_mu.get(bloque.codigo)
        if bloque_mu:
            _slide_volumen_multi(prs, bloque.codigo, periodo, ws_multi, bloque_mu)

    for titulo, bloques_marca in COBER_SLIDES:
        _slide_totales_cobertura(prs, titulo, bloques_marca, periodo, ws_cober)

    for bloque in bloques_av:
        bloque_cb = bloques_cb.get(bloque.codigo)
        if not bloque_cb:
            continue
        for titulo, bloques_marca in COBER_SLIDES:
            _slide_cobertura(prs, titulo, bloques_marca, bloque.codigo, periodo,
                             ws_cober, bloque_cb)

    salida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(salida)
    return len(prs.slides._sldIdLst)


def main() -> int:
    parser = argparse.ArgumentParser(description="PowerPoint del AVANCE BADIE")
    parser.add_argument("--archivo", required=True,
                        help="xlsx del avance (se lee, nunca se escribe)")
    parser.add_argument("--salida", help="ruta del pptx (por defecto, junto al xlsx)")
    parser.add_argument("--force", action="store_true",
                        help="sobrescribir el pptx si ya existe")
    parser.add_argument("--diagnostico", action="store_true",
                        help="listar las filas de total del libro que no cierran")
    args = parser.parse_args()

    archivo = Path(args.archivo)
    if not archivo.is_file():
        print(f"ERROR: no existe {archivo}", file=sys.stderr)
        return 1

    salida = Path(args.salida) if args.salida else archivo.with_suffix(".pptx")
    if salida.exists() and not args.force:
        print(f"ERROR: {salida} ya existe. Usar --force para sobrescribir.", file=sys.stderr)
        return 1

    slides = construir(archivo, salida, diagnostico=args.diagnostico)
    print(f"OK: {slides} slides -> {salida}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
