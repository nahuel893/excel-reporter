"""Build Marca.pptx and Generico.pptx decks via python-pptx.

Each deck has one slide per (zona, generico) combination — cobertura plus
optional comparacion. Slides are image-only (title textbox + full-slide
picture). Missing PNGs are skipped with a WARN log.
"""
from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.services.graficos_cobertura.constants import (
    PPTX_FONT_NAME,
    PPTX_GENERICO_FILENAME,
    PPTX_SLIDE_HEIGHT_IN,
    PPTX_SLIDE_WIDTH_IN,
    PPTX_TITLE_COLOR,
    ZONA_SLUGS,
)


logger = logging.getLogger(__name__)


_GENERICOS_MARCA = [
    ("CERVEZAS", "cervezas"),
    ("AGUAS SABORIZADAS", "aguas_saborizadas"),
    ("AGUAS MINERAL", "aguas_mineral"),
]

_GENERICOS_ALL = [
    ("CERVEZAS", "cervezas"),
    ("AGUAS SABORIZADAS", "aguas_saborizadas"),
    ("AGUAS MINERAL", "aguas_mineral"),
    ("SIDRAS Y LICORES", "sidras_y_licores"),
    ("VINOS CCU", "vinos_ccu"),
]

_DEFAULT_ZONAS = [
    ("NOA NORTE", ZONA_SLUGS["NOA NORTE"]),
    ("SALTA CAPITAL", ZONA_SLUGS["SALTA CAPITAL"]),
    ("INTERIOR SALTA SUR", ZONA_SLUGS["INTERIOR SALTA SUR"]),
    ("INTERIOR SALTA NORTE", ZONA_SLUGS["INTERIOR SALTA NORTE"]),
    ("JUJUY INTERIOR", ZONA_SLUGS["JUJUY INTERIOR"]),
]


def _set_font(run, size: int, bold: bool, color: RGBColor) -> None:
    run.font.name = PPTX_FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title(slide, prs, title_text: str) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.1), Inches(12), Inches(0.55),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    _set_font(run, size=28, bold=True, color=RGBColor(*PPTX_TITLE_COLOR))


def _add_picture(slide, prs, img_path: Path, margin_top: float) -> None:
    margin_x = Inches(0.3)
    top = Inches(margin_top)
    img_w = prs.slide_width - 2 * margin_x
    img_h = prs.slide_height - top - Inches(0.2)
    slide.shapes.add_picture(str(img_path), margin_x, top, img_w, img_h)


def _add_slide_with_image(prs, img_path: Path, title_text: str, margin_top: float) -> bool:
    """Add one image slide to the deck. Returns True if added, False if PNG missing."""
    if not img_path.exists():
        logger.warning("PPTX: missing png skipped: %s", img_path)
        return False

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_title(slide, prs, title_text)
    _add_picture(slide, prs, img_path, margin_top)
    return True


def build_pptx(
    genericos: list[tuple[str, str]],
    zonas: list[tuple[str, str]],
    png_dir: Path,
    output_path: Path,
    include_comparacion: bool = True,
) -> Path:
    """Build a single deck from PNG files in png_dir.

    For each (generico, zona) pair, adds:
      - one "cobertura" slide (always)
      - one "comparacion" slide (if include_comparacion=True)

    Missing PNGs are skipped with a WARN.
    """
    png_dir = Path(png_dir)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(PPTX_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(PPTX_SLIDE_HEIGHT_IN)

    for gen_name, gen_slug in genericos:
        # Primero TODAS las coberturas del genérico (5 zonas)
        for zona_name, zona_slug in zonas:
            cob_png = png_dir / f"cobertura_{zona_slug}_{gen_slug}.png"
            _add_slide_with_image(
                prs, cob_png,
                title_text=f"{zona_name} — {gen_name}",
                margin_top=0.95,
            )

        # Después TODAS las comparaciones del genérico (5 zonas)
        if include_comparacion:
            for zona_name, zona_slug in zonas:
                comp_png = png_dir / f"comparacion_{zona_slug}_{gen_slug}.png"
                _add_slide_with_image(
                    prs, comp_png,
                    title_text=f"{zona_name} — {gen_name} (Comparativo)",
                    margin_top=0.7,
                )

    prs.save(str(output_path))
    return output_path


def build_decks(
    png_dir: Path,
    output_dir: Path,
    con_aguas: bool = True,
    zonas: list[tuple[str, str]] | None = None,
) -> dict[str, Path]:
    """Build cobertura_todos.pptx from the PNG directory (all 5 genericos)."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    zonas = zonas or _DEFAULT_ZONAS

    if con_aguas:
        gens_all = _GENERICOS_ALL
    else:
        gens_all = [g for g in _GENERICOS_ALL if "AGUAS" not in g[0]]

    generico_path = build_pptx(
        genericos=gens_all,
        zonas=zonas,
        png_dir=png_dir,
        output_path=output_dir / PPTX_GENERICO_FILENAME,
    )

    return {"generico": generico_path}


def _resolve_zona_name(zona_slug: str) -> str:
    """Reverse-lookup zona slug to display name."""
    for name, slug in ZONA_SLUGS.items():
        if slug == zona_slug:
            return name
    return zona_slug.replace("_", " ").title()


def build_deck_sucursal(
    zona_slug: str,
    generico: str,
    sucursal_id: int,
    sucursal_nombre: str,
    png_paths_zone: dict[str, Path],
    png_paths_sucursal: dict[str, Path],
    output_path: Path,
) -> Path:
    """Build a per-sucursal PPTX with zone overview + sucursal detail slides.

    Args:
        zona_slug: Zone identifier (e.g. "salta_capital").
        generico: Generic name (e.g. "CERVEZAS").
        sucursal_id: Numeric sucursal ID.
        sucursal_nombre: Display name (e.g. "SUCURSAL METAN").
        png_paths_zone: Dict with keys "cobertura"/"comparacion" → Path for zone charts.
        png_paths_sucursal: Dict with keys "cobertura"/"comparacion" → Path for sucursal charts.
        output_path: Destination path for the PPTX file.

    Returns:
        Path to the created PPTX.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    zona_name = _resolve_zona_name(zona_slug)

    prs = Presentation()
    prs.slide_width = Inches(PPTX_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(PPTX_SLIDE_HEIGHT_IN)

    # Zone-level slides first
    for chart_type, margin_top in (("cobertura", 0.95), ("comparacion", 0.7)):
        img_path = png_paths_zone.get(chart_type)
        if img_path is None:
            continue
        suffix = " (Comparativo)" if chart_type == "comparacion" else ""
        _add_slide_with_image(
            prs, Path(img_path),
            title_text=f"{zona_name} — {generico}{suffix}",
            margin_top=margin_top,
        )

    # Sucursal-level detail slides
    for chart_type, margin_top in (("cobertura", 0.95), ("comparacion", 0.7)):
        img_path = png_paths_sucursal.get(chart_type)
        if img_path is None:
            continue
        suffix = " (Comparativo)" if chart_type == "comparacion" else ""
        _add_slide_with_image(
            prs, Path(img_path),
            title_text=f"{sucursal_nombre} — {zona_name} — {generico}{suffix}",
            margin_top=margin_top,
        )

    prs.save(str(output_path))
    return output_path
