"""Tests for `scripts/avance_branca_pptx.py`.

The fixture rebuilds the shape of the real workbook in miniature: categories on
the rows, three columns per vendor, a total row recognised by its label, and a
label-only row like the sheet's "(vinos fecovita)".

The number that matters here is the coverage total. It is NOT the sum of the
brand rows — a client that buys two brands is one point of sale — so the deck
has to read it from the sheet. The fixture makes them differ on purpose.
"""
import importlib.util
import sys
from pathlib import Path

import openpyxl
import pytest
from pptx import Presentation


RAIZ = Path(__file__).resolve().parents[1]
for ruta in (RAIZ, RAIZ / "scripts"):
    if str(ruta) not in sys.path:
        sys.path.insert(0, str(ruta))

_spec = importlib.util.spec_from_file_location(
    "avance_branca_pptx", RAIZ / "scripts" / "avance_branca_pptx.py")
branca = importlib.util.module_from_spec(_spec)
sys.modules["avance_branca_pptx"] = branca
_spec.loader.exec_module(branca)


VENDEDORES = ["PABLO NAVARRO", "GABRIEL JAUREGUI", "GONZALO LOPEZ", "CAROLINA DURAN",
              "LUCIANO GUZMAN", "NICOLAS NOVELLO", "FLORENCIA MEDINA",
              "SUB DISTRIBUIDOR", "NDIOLI", "GFARAH"]

# Coverage: what each brand row holds, and what the sheet's own total says.
COBER_FILAS = {"ANTICA": 0, "BORGHETTI": 24, "CARPANO": 101, "FERNET": 198}
COBER_TOTAL_DECLARADO = 201          # la hoja: PDV distintos
COBER_SUMA_DE_MARCAS = 323           # 0 + 24 + 101 + 198


def _set(ws, fila, columna, valor):
    ws.cell(fila, openpyxl.utils.column_index_from_string(columna)).value = valor


@pytest.fixture
def libro(tmp_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = branca.HOJA_AVANCE
    _set(ws, 1, "C", __import__("datetime").datetime(2026, 7, 31))
    _set(ws, 1, "I", 25)
    _set(ws, 2, "I", 25)
    for (columna, _cols), nombre in zip(branca.AVANCE_BLOQUES, VENDEDORES):
        _set(ws, branca.FILA_VENDEDORES, columna, nombre)
    _set(ws, 7, "A", "Categoria")

    # Volumen: dos marcas y el total. Avance / %Tend / Faltan por vendedor.
    for fila, etiqueta, avance, faltan in [(8, "ANTICA", 10.0, 2.0),
                                           (9, "FERNET", 30.0, 8.0)]:
        _set(ws, fila, "A", etiqueta)
        for _columna, cols in branca.AVANCE_BLOQUES:
            _set(ws, fila, cols[0], avance)
            _set(ws, fila, cols[1], avance / (avance + faltan))
            _set(ws, fila, cols[2], faltan)
    _set(ws, 22, "A", "TOTAL LINEA  Branca")
    for _columna, cols in branca.AVANCE_BLOQUES:
        _set(ws, 22, cols[0], 40.0)
        _set(ws, 22, cols[1], 0.8)
        _set(ws, 22, cols[2], 10.0)

    cob = wb.create_sheet(branca.HOJA_COBERTURA)
    for (columna, _cols), nombre in zip(branca.COBERTURA_BLOQUES, VENDEDORES):
        _set(cob, branca.FILA_VENDEDORES, columna, nombre)
    for fila, (etiqueta, valor) in enumerate(COBER_FILAS.items(), start=7):
        _set(cob, fila, "A", etiqueta)
        for _columna, cols in branca.COBERTURA_BLOQUES:
            _set(cob, fila, cols[0], valor)
    _set(cob, 20, "A", "(vinos fecovita)")   # fila de aclaracion, sin datos
    _set(cob, 21, "A", "TOTAL LINEA")
    for _columna, cols in branca.COBERTURA_BLOQUES:
        _set(cob, 21, cols[0], COBER_TOTAL_DECLARADO)

    ruta = tmp_path / "AVANCE BRANCA - JULIO 2026.xlsx"
    wb.save(ruta)
    return ruta


def _tablas(pptx_path):
    """[(titulo, [[texto de celda, ...], ...]), ...] por slide con tabla."""
    salida = []
    for slide in Presentation(pptx_path).slides:
        titulo = ""
        tabla = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() and not titulo:
                titulo = shape.text_frame.text.split("\n")[0]
            if shape.has_table:
                tabla = shape.table
        if tabla is not None:
            filas = [[tabla.cell(r, c).text for c in range(len(tabla.columns))]
                     for r in range(len(tabla.rows))]
            salida.append((titulo, filas))
    return salida


@pytest.fixture
def deck(libro, tmp_path):
    salida = tmp_path / "deck.pptx"
    branca.construir(libro, salida, rechazos=None, con_capturas=False)
    return _tablas(salida)


class TestLecturaDeLaHoja:
    def test_reconoce_que_el_libro_trae_valores(self, libro):
        wb = openpyxl.load_workbook(libro, data_only=True)
        assert branca._tiene_valores(wb[branca.HOJA_AVANCE]) is True

    def test_un_libro_sin_valores_cacheados_se_detecta(self, tmp_path):
        """`AvancesService` escribe con openpyxl y pierde los valores cacheados:
        sin esto el deck saldria lleno de guiones."""
        wb = openpyxl.Workbook()
        wb.active.title = branca.HOJA_AVANCE
        ruta = tmp_path / "vacio.xlsx"
        wb.save(ruta)
        recargado = openpyxl.load_workbook(ruta, data_only=True)
        assert branca._tiene_valores(recargado[branca.HOJA_AVANCE]) is False

    def test_las_etiquetas_de_vendedor_salen_de_la_fila_5(self, libro):
        wb = openpyxl.load_workbook(libro, data_only=True)
        etiquetas = branca._etiquetas_bloques(wb[branca.HOJA_AVANCE], branca.AVANCE_BLOQUES)
        assert etiquetas == VENDEDORES

    def test_la_fila_de_total_se_reconoce_por_su_etiqueta(self, libro):
        wb = openpyxl.load_workbook(libro, data_only=True)
        columnas = [c for _e, cols in branca.AVANCE_BLOQUES for c in cols]
        detalle, total = branca._filas_seccion(wb[branca.HOJA_AVANCE], columnas, 8, 22)
        assert [e for _f, e in detalle] == ["ANTICA", "FERNET"]
        assert total == (22, "TOTAL LINEA  Branca")

    def test_las_filas_de_solo_aclaracion_se_descartan(self, libro):
        """La hoja tiene renglones con texto en A y ninguna cifra."""
        wb = openpyxl.load_workbook(libro, data_only=True)
        columnas = [c for _e, cols in branca.COBERTURA_BLOQUES for c in cols]
        detalle, _total = branca._filas_seccion(wb[branca.HOJA_COBERTURA], columnas, 7, 21)
        assert "(vinos fecovita)" not in [e for _f, e in detalle]


class TestCoberturaNoSeSuma:
    """El total de cobertura son PDV distintos, no la suma de las marcas."""

    def _fila_total(self, deck, titulo_contiene):
        for titulo, filas in deck:
            if titulo_contiene in titulo:
                return filas[-1]
        raise AssertionError(f"no hay slide con '{titulo_contiene}'")

    def test_el_total_de_cobertura_es_el_de_la_hoja(self, deck):
        fila = self._fila_total(deck, "COBERTURA")
        assert fila[0] == "TOTAL LINEA"
        assert fila[1] == str(COBER_TOTAL_DECLARADO)

    def test_el_total_de_cobertura_no_es_la_suma_de_las_marcas(self, deck):
        fila = self._fila_total(deck, "COBERTURA")
        assert fila[1] != str(COBER_SUMA_DE_MARCAS)

    def test_el_total_de_volumen_si_es_la_suma(self, deck):
        """Los bultos si son aditivos entre categorias."""
        fila = self._fila_total(deck, "VOLUMEN")
        assert fila[0].startswith("TOTAL LINEA")
        assert fila[1] == "40"          # 10 + 30
        assert fila[3] == "10,0"        # 2 + 8; Faltan lleva un decimal, como la hoja

    def test_el_porcentaje_del_total_se_deriva_no_se_suma(self, deck):
        fila = self._fila_total(deck, "VOLUMEN")
        assert fila[2] == "80,0%"       # 40 / (40 + 10)


class TestOrdenDeLasDiapositivas:
    def test_primero_volumen_y_despues_cobertura(self, deck):
        titulos = [t for t, _f in deck]
        volumen = [i for i, t in enumerate(titulos) if t.startswith("VOLUMEN")]
        cobertura = [i for i, t in enumerate(titulos) if t.startswith("COBERTURA")]
        assert volumen and cobertura
        assert max(volumen) < min(cobertura)

    def test_hay_una_columna_por_vendedor_en_cobertura(self, deck):
        for titulo, filas in deck:
            if titulo.startswith("COBERTURA"):
                assert len(filas[0]) == 1 + len(branca.COBERTURA_BLOQUES)

    def test_hay_tres_columnas_por_vendedor_en_volumen(self, deck):
        for titulo, filas in deck:
            if titulo.startswith("VOLUMEN"):
                assert len(filas[0]) == 1 + 3 * len(branca.AVANCE_BLOQUES)


class TestRechazos:
    def test_la_imagen_se_busca_por_carpeta_y_no_por_nombre(self, tmp_path, monkeypatch):
        """`configs/rebotes.json` tiene el nombre escrito a mano, asi que todos
        los meses salen como "Rebotes Junio 2026". La carpeta si es correcta."""
        carpeta = tmp_path / "data" / "output" / "reporte-rebotes" / "2026-07"
        carpeta.mkdir(parents=True)
        png = carpeta / "Rebotes Junio 2026_% Rebotes x Generico_A1_K40.png"
        png.write_bytes(b"x")
        monkeypatch.setattr(branca, "RAIZ", tmp_path)
        assert branca._buscar_rechazos("2026-07") == png

    def test_sin_carpeta_del_mes_no_rompe(self, tmp_path, monkeypatch):
        monkeypatch.setattr(branca, "RAIZ", tmp_path)
        assert branca._buscar_rechazos("2026-07") is None

    def test_el_deck_sale_igual_sin_imagen(self, deck):
        assert any(t.startswith("VOLUMEN") for t, _f in deck)


class TestCapturasDeLaHoja:
    def test_las_capturas_de_backup_se_descartan(self, tmp_path):
        """Los `backup-*` son de una corrida vieja del mismo mes."""
        xlsx = tmp_path / "AVANCE BRANCA - JULIO 2026.xlsx"
        xlsx.touch()
        vigente = tmp_path / "AVANCE BRANCA - JULIO 2026_AVANCE_B2_AX35.png"
        vieja = tmp_path / "AVANCE BRANCA - JULIO 2026_AVANCE_B2_AX35_backup-20260708.png"
        for png in (vigente, vieja):
            png.write_bytes(b"x")
        assert branca._capturas(xlsx) == {"AVANCE": vigente}

    def test_reconoce_la_captura_de_cada_hoja(self, tmp_path):
        xlsx = tmp_path / "AVANCE BRANCA - JULIO 2026.xlsx"
        xlsx.touch()
        avance = tmp_path / "AVANCE BRANCA - JULIO 2026_AVANCE_B2_AX35.png"
        cober = tmp_path / "AVANCE BRANCA - JULIO 2026_Cobertura_B2_AR37.png"
        for png in (avance, cober):
            png.write_bytes(b"x")
        assert branca._capturas(xlsx) == {"AVANCE": avance, "COBERTURA": cober}
