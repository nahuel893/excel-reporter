"""Tests for `scripts/avance_deck.py` — the single monthly deck.

The two builders (`avance_pptx`, `avance_branca_pptx`) already have their own
suites; this module does not re-test them. What it fixes is the part only the
orchestrator owns:

- the order of the sections, because that order is the meeting's running order;
- that RECHAZOS appears ONCE, at deck level. The BRANCA builder puts it at the
  end of its own deck, so wiring it in naively would print it twice;
- that a missing image drops its slide instead of failing the whole deck.
"""
import importlib.util
import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches


RAIZ = Path(__file__).resolve().parents[1]
for ruta in (RAIZ, RAIZ / "scripts"):
    if str(ruta) not in sys.path:
        sys.path.insert(0, str(ruta))

_spec = importlib.util.spec_from_file_location(
    "avance_deck", RAIZ / "scripts" / "avance_deck.py")
deck = importlib.util.module_from_spec(_spec)
sys.modules["avance_deck"] = deck
_spec.loader.exec_module(deck)


def _png(destino: Path, ancho: int = 400, alto: int = 300) -> Path:
    from PIL import Image

    Image.new("RGB", (ancho, alto), (200, 210, 230)).save(destino)
    return destino


def _titulos(pptx: Path) -> list:
    """El primer texto de cada slide, que es siempre el titulo."""
    titulos = []
    for slide in Presentation(str(pptx)).slides:
        for forma in slide.shapes:
            if forma.has_text_frame and forma.text_frame.text.strip():
                titulos.append(forma.text_frame.text.strip().splitlines()[0])
                break
        else:
            titulos.append("")
    return titulos


@pytest.fixture
def builders(monkeypatch):
    """Reemplaza los dos builders por marcadores; deja el orden a la vista."""
    llamadas = []

    def _marca(texto):
        def _poblar(prs, archivo, **kwargs):
            llamadas.append((texto, archivo, kwargs))
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            caja = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            caja.text_frame.text = texto
        return _poblar

    monkeypatch.setattr(deck.badie, "poblar", _marca("BADIE"))
    monkeypatch.setattr(deck.branca, "poblar", _marca("BRANCA"))
    monkeypatch.setattr(deck, "_periodo", lambda archivo: "JULIO 2026")
    return llamadas


class TestOrdenDeLasSecciones:
    def test_badie_branca_rechazos_danielito(self, tmp_path, builders):
        salida = tmp_path / "deck.pptx"
        deck.construir(
            tmp_path / "badie.xlsx", salida,
            archivo_branca=tmp_path / "branca.xlsx",
            rechazos=_png(tmp_path / "rechazos.png"),
            danielito=_png(tmp_path / "danielito.png"),
        )
        assert _titulos(salida) == ["BADIE", "BRANCA", "RECHAZOS", "VINOS DANIELITO"]

    def test_sin_branca_el_deck_sigue_saliendo(self, tmp_path, builders):
        salida = tmp_path / "deck.pptx"
        deck.construir(tmp_path / "badie.xlsx", salida,
                       rechazos=_png(tmp_path / "rechazos.png"))
        assert _titulos(salida) == ["BADIE", "RECHAZOS"]


class TestRechazosUnaSolaVez:
    def test_el_builder_de_branca_no_pone_su_propia_slide(self, tmp_path, builders):
        """El deck unificado pone RECHAZOS al final, no adentro de BRANCA."""
        deck.construir(tmp_path / "badie.xlsx", tmp_path / "deck.pptx",
                       archivo_branca=tmp_path / "branca.xlsx",
                       rechazos=_png(tmp_path / "rechazos.png"))
        _texto, _archivo, kwargs = builders[1]
        assert kwargs["con_rechazos"] is False

    def test_rechazos_aparece_una_sola_vez(self, tmp_path, builders):
        salida = tmp_path / "deck.pptx"
        deck.construir(tmp_path / "badie.xlsx", salida,
                       archivo_branca=tmp_path / "branca.xlsx",
                       rechazos=_png(tmp_path / "rechazos.png"))
        assert _titulos(salida).count("RECHAZOS") == 1


class TestImagenesQueFaltan:
    def test_sin_rechazos_no_hay_slide_de_rechazos(self, tmp_path, builders):
        salida = tmp_path / "deck.pptx"
        deck.construir(tmp_path / "badie.xlsx", salida,
                       danielito=_png(tmp_path / "danielito.png"))
        assert "RECHAZOS" not in _titulos(salida)

    def test_un_png_inexistente_no_rompe_el_deck(self, tmp_path, builders):
        salida = tmp_path / "deck.pptx"
        deck.construir(tmp_path / "badie.xlsx", salida,
                       rechazos=tmp_path / "no-existe.png")
        assert _titulos(salida) == ["BADIE"]


class TestBusquedaDeArchivos:
    def test_branca_se_busca_al_lado_del_libro_de_badie(self, tmp_path):
        (tmp_path / "AVANCE BADIE - JULIO 2026.xlsx").touch()
        esperado = tmp_path / "AVANCE BRANCA - JULIO 2026.xlsx"
        esperado.touch()
        assert deck._buscar_branca(tmp_path / "AVANCE BADIE - JULIO 2026.xlsx") == esperado

    def test_sin_libro_de_branca_devuelve_none(self, tmp_path):
        badie = tmp_path / "AVANCE BADIE - JULIO 2026.xlsx"
        badie.touch()
        assert deck._buscar_branca(badie) is None

    def test_danielito_se_busca_por_carpeta_de_periodo(self, tmp_path, monkeypatch):
        """El nombre del PNG cambia con el rango capturado; la carpeta no."""
        carpeta = tmp_path / "data" / "output" / "vinos-danielito" / "2026-07"
        carpeta.mkdir(parents=True)
        _png(carpeta / "Vinos Danielito - Volumen por mes_Volumen por mes_A1_O28.png")
        monkeypatch.setattr(deck, "RAIZ", tmp_path)
        assert deck._buscar_danielito("2026-07").parent == carpeta

    def test_danielito_descarta_los_backup(self, tmp_path, monkeypatch):
        """Un `backup-*` puede tener fecha de archivo mas nueva que el vigente."""
        carpeta = tmp_path / "data" / "output" / "vinos-danielito" / "2026-07"
        carpeta.mkdir(parents=True)
        vigente = _png(carpeta / "Vinos Danielito - Volumen por mes_A1_O28.png")
        backup = _png(carpeta / "Vinos Danielito - Volumen por mes_A1_O28_backup-20260822.png")
        import os
        os.utime(backup, (backup.stat().st_atime + 1000, backup.stat().st_mtime + 1000))
        monkeypatch.setattr(deck, "RAIZ", tmp_path)
        assert deck._buscar_danielito("2026-07") == vigente

    def test_danielito_sin_carpeta_devuelve_none(self, tmp_path, monkeypatch):
        monkeypatch.setattr(deck, "RAIZ", tmp_path)
        assert deck._buscar_danielito("2026-07") is None
