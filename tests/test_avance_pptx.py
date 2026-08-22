"""Tests for `scripts/avance_pptx.py`.

The fixtures rebuild the shape of the real workbook in miniature: two supervisor
bands, a roll-up band whose members are supervisor codes, and the two
single-line bands (DIRECTA, SUB DISTRIBUIDOR). That shape is what the block
detection has to survive.
"""
import importlib.util
import sys
from pathlib import Path

import openpyxl
import pytest


RAIZ = Path(__file__).resolve().parents[1]
_spec = importlib.util.spec_from_file_location("avance_pptx", RAIZ / "scripts" / "avance_pptx.py")
avance_pptx = importlib.util.module_from_spec(_spec)
# @dataclass resolves annotations through sys.modules, so register before exec.
sys.modules["avance_pptx"] = avance_pptx
_spec.loader.exec_module(avance_pptx)


@pytest.fixture
def hoja_avance():
    """`Avance`-shaped sheet: B=vendedor, C=supervisor."""
    wb = openpyxl.Workbook()
    ws = wb.active
    filas = [
        (7, "Vendedor", "Super"),          # header, must be ignored
        (8, "ANA PEREZ", "GFLORES"),
        (9, "LUIS GOMEZ", "GFLORES"),
        (11, "MARIA RUIZ", "FGUANTAY"),
        (12, "JUAN DIAZ", "FGUANTAY"),
        (14, "DIRECTA", "GFARAH"),         # single-line band
        (15, "SUB DISTRIBUIDOR", "ANOGALES"),
        (17, "GFLORES", "GFARAH"),         # summary band: members are codes
        (18, "FGUANTAY", "GFARAH"),
    ]
    for fila, nombre, codigo in filas:
        ws.cell(fila, 2).value = nombre
        ws.cell(fila, 3).value = codigo
    return ws


class TestDeteccionDeBloques:
    def test_solo_los_supervisores_son_bloques(self, hoja_avance):
        codigos = [b.codigo for b in avance_pptx._bloques(hoja_avance, "B", "C", 20)]
        assert codigos == ["GFLORES", "FGUANTAY"]

    def test_la_banda_de_resumen_no_es_un_bloque(self, hoja_avance):
        """GFARAH agrupa codigos, no vendedores: no puede tener slide propio."""
        codigos = [b.codigo for b in avance_pptx._bloques(hoja_avance, "B", "C", 20)]
        assert "GFARAH" not in codigos

    def test_las_bandas_de_una_sola_linea_no_son_bloques(self, hoja_avance):
        codigos = [b.codigo for b in avance_pptx._bloques(hoja_avance, "B", "C", 20)]
        assert "ANOGALES" not in codigos

    def test_el_bloque_conserva_el_orden_de_la_hoja(self, hoja_avance):
        bloque = avance_pptx._bloques(hoja_avance, "B", "C", 20)[0]
        assert bloque.vendedores == ["ANA PEREZ", "LUIS GOMEZ"]
        assert bloque.filas == {"ANA PEREZ": 8, "LUIS GOMEZ": 9}

    def test_el_encabezado_no_entra_como_vendedor(self, hoja_avance):
        bloques = avance_pptx._bloques(hoja_avance, "B", "C", 20)
        assert all("Vendedor" not in b.vendedores for b in bloques)

    def test_rollup_es_el_codigo_que_agrupa_codigos(self, hoja_avance):
        assert avance_pptx._codigo_rollup(hoja_avance, "B", "C", 20) == "GFARAH"

    def test_otros_son_directa_y_sub_distribuidor(self, hoja_avance):
        sueltos = avance_pptx._otros(hoja_avance, "B", "C", 20, {"GFLORES", "FGUANTAY"})
        assert sueltos == [("DIRECTA", 14), ("SUB DISTRIBUIDOR", 15)]


class TestSumas:
    """Totals are the sum of the rows shown, never a cell of the workbook.

    Two of the workbook's own total rows are stale (see the module docstring),
    so reading them would put a total on the slide that its own rows contradict.
    """

    @pytest.fixture
    def hoja(self):
        wb = openpyxl.Workbook()
        ws = wb.active
        for fila, venta in [(2, 100.5), (3, 200.25), (4, None)]:
            ws.cell(fila, 4).value = venta
        ws.cell(5, 4).value = "#REF!"
        return ws

    def test_suma_ignora_vacios(self, hoja):
        assert avance_pptx._sumar(hoja, [2, 3, 4], "D") == pytest.approx(300.75)

    def test_suma_ignora_errores_de_formula(self, hoja):
        assert avance_pptx._sumar(hoja, [2, 3, 5], "D") == pytest.approx(300.75)

    def test_suma_sin_ningun_numero_es_none(self, hoja):
        assert avance_pptx._sumar(hoja, [4, 5], "D") is None

    def test_la_suma_no_redondea(self, hoja):
        ws = hoja
        ws.cell(6, 4).value = 0.3333333333
        assert avance_pptx._sumar(ws, [6], "D") == 0.3333333333


class TestRatio:
    def test_ratio_normal(self):
        assert avance_pptx._ratio(90, 100) == pytest.approx(0.9)

    def test_denominador_cero_no_rompe(self):
        """DIRECTA no tiene cupo: la celda queda en '-', no en 0 ni en #DIV/0!."""
        assert avance_pptx._ratio(2054.18, 0) is None

    def test_numerador_none(self):
        assert avance_pptx._ratio(None, 100) is None


class TestFormato:
    def test_miles_con_punto_y_venta_sin_decimales(self):
        assert avance_pptx._fmt(23340.65, "num") == "23.341"

    def test_pdv_sin_decimales(self):
        assert avance_pptx._fmt(1593, "pdv") == "1.593"

    def test_el_objetivo_de_cobertura_conserva_un_decimal(self):
        """En SIDRAS el objetivo anda por 9,3: redondearlo lo deja sin sentido."""
        assert avance_pptx._fmt(9.3116883, "obj") == "9,3"

    def test_porcentaje_con_un_decimal(self):
        assert avance_pptx._fmt(0.9695, "pct") == "97,0%"
        assert avance_pptx._fmt(0.9644, "pct") == "96,4%"

    def test_un_porcentaje_disparado_pierde_el_decimal(self):
        """AMSTEL con objetivo 2,3 da 1.027,8%: el decimal parte la celda en dos."""
        assert avance_pptx._fmt(10.278, "pct") == "1.028%"

    def test_vacio_se_muestra_como_guion(self):
        assert avance_pptx._fmt(None, "num") == "-"

    def test_los_errores_de_formula_no_son_numero(self):
        assert avance_pptx._numero("#DIV/0!") is None
        assert avance_pptx._numero(True) is None
        assert avance_pptx._numero(0) == 0


class TestSemaforo:
    def test_cumplido_va_en_verde(self):
        assert avance_pptx._color_pct(1.04) == avance_pptx.VERDE

    def test_entre_90_y_100_va_en_ambar(self):
        assert avance_pptx._color_pct(0.93) == avance_pptx.AMBAR

    def test_debajo_de_90_va_en_rojo(self):
        assert avance_pptx._color_pct(0.72) == avance_pptx.ROJO

    def test_sin_dato_no_pinta(self):
        assert avance_pptx._color_pct(None) is None


class TestColumnasAditivas:
    """A percentage is a ratio: summing it across rows is meaningless."""

    def test_las_columnas_de_porcentaje_quedan_fuera(self):
        aditivas = set(avance_pptx._cols_aditivas_avance())
        porcentajes = {pct for _n, _v, pct, _f in avance_pptx.AVANCE_CATEGORIAS}
        porcentajes |= {col for _e, col, clase in avance_pptx.AVANCE_TOTAL_COLS
                        if clase == "pct"}
        assert not (aditivas & porcentajes)

    def test_cobertura_suma_todo_menos_el_porcentaje(self):
        aditivas = set(avance_pptx._cols_aditivas_cobertura())
        for _titulo, bloques in avance_pptx.COBER_SLIDES:
            for _marca, columnas in bloques:
                assert set(columnas[:-1]) <= aditivas
                assert columnas[-1] not in aditivas


class TestDiagnostico:
    def test_detecta_una_fila_de_total_desactualizada(self):
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.cell(2, 4).value = 100.0
        ws.cell(3, 4).value = 50.0
        ws.cell(4, 4).value = 100.0  # total row that forgot row 3
        bloque = avance_pptx.Bloque(codigo="GFLORES",
                                    vendedores=["A", "B"],
                                    filas={"A": 2, "B": 3})
        avisos = avance_pptx._diagnostico(ws, "Avance", bloque, ["D"], 4, "total del bloque")
        assert len(avisos) == 1
        assert "150.00" in avisos[0] and "100.00" in avisos[0]

    def test_no_avisa_cuando_el_total_cierra(self):
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.cell(2, 4).value = 100.0
        ws.cell(3, 4).value = 50.0
        ws.cell(4, 4).value = 150.0
        bloque = avance_pptx.Bloque(codigo="GFLORES",
                                    vendedores=["A", "B"],
                                    filas={"A": 2, "B": 3})
        assert avance_pptx._diagnostico(ws, "Avance", bloque, ["D"], 4, "total") == []


class TestLayoutDeColumnas:
    """La hoja da tres columnas por categoria y cinco en el total: van todas."""

    def test_cada_categoria_de_cervezas_trae_sus_tres_columnas(self):
        cols = avance_pptx._cols_volumen_avance()
        assert len(cols) == len(avance_pptx.AVANCE_CATEGORIAS) * 3 + len(avance_pptx.AVANCE_TOTAL_COLS)
        for _nombre, venta, pct, falta in avance_pptx.AVANCE_CATEGORIAS:
            assert venta in cols and pct in cols and falta in cols

    def test_el_total_cerveza_trae_cinco_columnas(self):
        assert [e for e, _c, _cl in avance_pptx.AVANCE_TOTAL_COLS] == [
            "Venta", "Cupo", "Tend.", "% Tend.", "Dif."]

    def test_la_venta_diaria_por_cupo_queda_fuera(self):
        """AR es #DIV/0! en todas las filas del libro: mostrarla es una columna de guiones."""
        assert "AR" not in [c for _e, c, _cl in avance_pptx.AVANCE_TOTAL_COLS]

    def test_multicategoria_tambien_trae_tres_columnas(self):
        cols = avance_pptx._cols_volumen_multi()
        assert len(cols) == len(avance_pptx.MULTI_CATEGORIAS) * 3

    def test_cobertura_se_abre_por_marca_no_solo_por_generico(self):
        marcas = [m for _t, bloques in avance_pptx.COBER_SLIDES for m, _c in bloques]
        for marca in ["SALTA", "SCHNEIDER", "LEVITE", "VILLA DEL SUR", "COLON", "MISTRAL"]:
            assert marca in marcas

    def test_cada_bloque_de_cobertura_cierra_con_su_total(self):
        for _titulo, bloques in avance_pptx.COBER_SLIDES:
            if _titulo == "COBERTURA CERVEZAS 1":
                continue  # su total vive en el slide de Cervezas 2
            assert bloques[-1][0].startswith("TOTAL")

    def test_las_marcas_de_vinos_y_sidras_tienen_tres_columnas(self):
        """La hoja no les calcula 'Faltan'; inventarlo seria dato nuevo."""
        for titulo, bloques in avance_pptx.COBER_SLIDES:
            if "VINOS" not in titulo and "SIDRAS" not in titulo:
                continue
            for marca, columnas in bloques:
                assert len(columnas) == (4 if marca.startswith("TOTAL") else 3)

    def test_los_cuatro_genericos_pedidos_estan_y_pernod_no(self):
        nombres = [n for n, *_ in avance_pptx.COBER_GENERICOS]
        assert nombres == ["CERVEZAS", "AGUAS DANONE", "VINOS CCU", "SIDRAS Y LICORES"]
        titulos = " ".join(t for t, _b in avance_pptx.COBER_SLIDES)
        assert "PERNOD" not in titulos

    def test_ninguna_columna_se_repite_dentro_de_un_slide(self):
        for titulo, bloques in avance_pptx.COBER_SLIDES:
            columnas = [c for _m, cols in bloques for c in cols]
            assert len(columnas) == len(set(columnas)), titulo


class TestTotalesPorSupervisor:
    """La slide que abre cada seccion: una fila por supervisor, sus totales.

    Antes habia una sola slide de resumen que mezclaba volumen (azul) y
    cobertura (verde) en la misma tabla. Son dos lecturas distintas y juntarlas
    obliga a saltar de una a otra; ahora cada seccion arranca con los totales de
    su propia hoja.

    Estas filas se LEEN de la banda de resumen del libro, celda por celda. No se
    suman los vendedores ni se recalcula ningun porcentaje: el numero que se
    muestra es el que Nahuel ve en el Excel. Si la hoja y la suma de sus filas no
    coinciden, eso lo reporta `--diagnostico`, no lo arregla el deck por su
    cuenta.
    """

    @pytest.fixture
    def hoja(self):
        """Vendedores arriba, banda de resumen abajo, rollup al pie.

        Los valores de la banda NO son la suma de los vendedores a proposito:
        asi un test puede distinguir "lei la celda" de "sume las filas".
        """
        wb = openpyxl.Workbook()
        ws = wb.active
        vendedores = [
            (8, "ANA PEREZ", "GFLORES", 45, 5),
            (9, "LUIS GOMEZ", "GFLORES", 45, 5),
            (11, "MARIA RUIZ", "FGUANTAY", 50, 450),
            (12, "JUAN DIAZ", "FGUANTAY", 50, 450),
        ]
        for fila, nombre, codigo, venta, falta in vendedores:
            ws.cell(fila, 2).value = nombre
            ws.cell(fila, 3).value = codigo
            ws.cell(fila, 4).value = venta
            ws.cell(fila, 6).value = falta

        # Banda de resumen: nombre = codigo del supervisor, super = rollup.
        for fila, codigo, venta, falta in [(17, "GFLORES", 111, 22),
                                           (18, "FGUANTAY", 333, 44)]:
            ws.cell(fila, 2).value = codigo
            ws.cell(fila, 3).value = "GFARAH"
            ws.cell(fila, 4).value = venta
            ws.cell(fila, 6).value = falta
        return ws

    def _filas(self, hoja):
        return avance_pptx._filas_por_supervisor(
            hoja, "B", "C", 20, ["D", "F"])

    def test_una_fila_por_supervisor_en_el_orden_de_la_hoja(self, hoja):
        etiquetas = [etiqueta for etiqueta, _valores, _total in self._filas(hoja)]

        assert etiquetas == ["GFLORES", "FGUANTAY"]

    def test_toma_el_dato_de_la_banda_y_no_suma_los_vendedores(self, hoja):
        """GFLORES: la banda dice 111, sus dos vendedores suman 90."""
        valores = dict((e, v) for e, v, _t in self._filas(hoja))

        assert valores["GFLORES"] == [111, 22]
        assert valores["FGUANTAY"] == [333, 44]

    def test_no_toca_el_valor_leido(self, hoja):
        """Ni redondeo ni escala: la celda entra al deck como esta."""
        hoja.cell(17, 4).value = 1234.567

        valores = dict((e, v) for e, v, _t in self._filas(hoja))

        assert valores["GFLORES"][0] == 1234.567

    def test_una_celda_vacia_de_la_banda_queda_en_none(self, hoja):
        """Sin dato no se inventa un cero: la celda del deck queda en '-'."""
        hoja.cell(18, 6).value = None

        valores = dict((e, v) for e, v, _t in self._filas(hoja))

        assert valores["FGUANTAY"][1] is None

    def test_un_error_de_formula_no_entra_como_dato(self, hoja):
        hoja.cell(17, 4).value = "#DIV/0!"

        valores = dict((e, v) for e, v, _t in self._filas(hoja))

        assert valores["GFLORES"][0] is None

    def test_ningun_supervisor_va_marcado_como_fila_de_total(self, hoja):
        assert all(total is False for _e, _v, total in self._filas(hoja))

    def test_sin_banda_de_resumen_no_hay_filas(self):
        """Si el libro no trae la banda, la slide no inventa una."""
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.cell(8, 2).value = "ANA PEREZ"
        ws.cell(8, 3).value = "GFLORES"
        ws.cell(9, 2).value = "LUIS GOMEZ"
        ws.cell(9, 3).value = "GFLORES"

        assert avance_pptx._filas_por_supervisor(ws, "B", "C", 20, ["D"]) == []

    def test_la_banda_de_cober_nueva_trae_un_cero_y_no_un_vacio(self):
        """`Cober Nueva` no tiene la misma forma que `Avance`.

        Ahi la banda de resumen (filas 51-54 del libro real) lleva un 0 numerico
        en la columna A en vez de dejarla vacia, y las filas de vendedor tambien
        llevan el codigo del supervisor en B. Buscar "filas cuyo nombre es un
        codigo" perdia la banda entera y la seccion de cobertura salia sin su
        slide de totales.
        """
        wb = openpyxl.Workbook()
        ws = wb.active
        for fila, nombre, codigo, pdv in [(8, "ANA PEREZ", "GFLORES", 10),
                                          (9, "LUIS GOMEZ", "GFLORES", 20)]:
            ws.cell(fila, 1).value = nombre      # A = vendedor
            ws.cell(fila, 2).value = codigo      # B = supervisor
            ws.cell(fila, 3).value = pdv
        ws.cell(51, 1).value = 0                 # la banda: A trae un 0
        ws.cell(51, 2).value = "GFLORES"
        ws.cell(51, 3).value = 1577

        filas = avance_pptx._filas_por_supervisor(
            ws, "A", "B", 55, ["C"], exigir_nombre_vacio=True)

        assert filas == [("GFLORES", [1577], False)]

    def test_la_ultima_fila_de_vendedor_no_se_confunde_con_la_banda(self):
        """Sin el flag, la fila 9 (B=GFLORES) ganaria por ser la ultima."""
        wb = openpyxl.Workbook()
        ws = wb.active
        for fila, nombre, pdv in [(8, "ANA PEREZ", 10), (9, "LUIS GOMEZ", 20)]:
            ws.cell(fila, 1).value = nombre
            ws.cell(fila, 2).value = "GFLORES"
            ws.cell(fila, 3).value = pdv
        ws.cell(51, 1).value = 0
        ws.cell(51, 2).value = "GFLORES"
        ws.cell(51, 3).value = 1577

        filas = avance_pptx._filas_por_supervisor(
            ws, "A", "B", 55, ["C"], exigir_nombre_vacio=True)

        assert filas[0][1] == [1577]   # la banda, no el 20 de LUIS GOMEZ


class TestBordesDeCelda:
    """Toda celda de la tabla lleva sus cuatro bordes.

    python-pptx no expone API de bordes: hay que meter los `a:lnL/R/T/B` a mano
    en el XML. `a:tcPr` tiene un esquema ORDENADO y las lineas van ANTES del
    relleno; si quedan despues de `a:solidFill`, PowerPoint se niega a abrir el
    archivo. Por eso el test de orden no es cosmetico.
    """

    @pytest.fixture
    def tabla(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(avance_pptx.ANCHO_IN)
        prs.slide_height = Inches(avance_pptx.ALTO_IN)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        spec = avance_pptx.Tabla(
            grupos=[("SALTA", ["Venta", "% Cupo"], avance_pptx.AZUL_OSCURO)],
            etiqueta_primera="Vendedor",
            filas=[("ANA PEREZ", [100, 0.9], False),
                   ("TOTAL", [100, 0.9], True)],
            clases=["num", "pct"],
        )
        avance_pptx._dibujar(slide, spec, top_in=1.05, ancho_primera=1.6)
        return next(sh.table for sh in slide.shapes if sh.has_table)

    def _lineas(self, celda):
        from pptx.oxml.ns import qn
        tc_pr = celda._tc.find(qn("a:tcPr"))
        if tc_pr is None:
            return []
        return [e.tag.split("}")[1] for e in tc_pr
                if e.tag.split("}")[1].startswith("ln")]

    def test_toda_celda_lleva_los_cuatro_bordes(self, tabla):
        faltan = []
        for f in range(len(tabla.rows)):
            for c in range(len(tabla.columns)):
                lineas = set(self._lineas(tabla.cell(f, c)))
                if not {"lnL", "lnR", "lnT", "lnB"} <= lineas:
                    faltan.append((f, c, sorted(lineas)))

        assert faltan == []

    def test_los_bordes_van_antes_del_relleno(self, tabla):
        """El orden del esquema: si `a:solidFill` queda primero, no abre."""
        from pptx.oxml.ns import qn

        for f in range(len(tabla.rows)):
            for c in range(len(tabla.columns)):
                tc_pr = tabla.cell(f, c)._tc.find(qn("a:tcPr"))
                hijos = [e.tag.split("}")[1] for e in tc_pr]
                relleno = [i for i, t in enumerate(hijos)
                           if t.endswith("Fill")]
                lineas = [i for i, t in enumerate(hijos) if t.startswith("ln")]
                if relleno and lineas:
                    assert max(lineas) < min(relleno), f"celda ({f},{c}): {hijos}"

    def test_el_orden_entre_lineas_es_el_del_esquema(self, tabla):
        """lnL, lnR, lnT, lnB — en ese orden, no cualquiera."""
        lineas = self._lineas(tabla.cell(2, 1))

        assert lineas == ["lnL", "lnR", "lnT", "lnB"]

    def test_el_borde_tiene_color_y_ancho(self, tabla):
        from pptx.oxml.ns import qn

        tc_pr = tabla.cell(2, 1)._tc.find(qn("a:tcPr"))
        ln = tc_pr.find(qn("a:lnL"))
        color = ln.find(qn("a:solidFill")).find(qn("a:srgbClr"))

        assert ln.get("w") == str(int(avance_pptx.BORDE_ANCHO_PT * 12700))
        assert color.get("val") == str(avance_pptx.BORDE_COLOR)

    def test_bordear_dos_veces_no_duplica(self):
        """Idempotente: duplicar las lineas deja el `a:tcPr` invalido."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        forma = slide.shapes.add_table(2, 2, 0, 0, 1000000, 1000000)
        celda = forma.table.cell(0, 0)

        avance_pptx._borde(celda)
        avance_pptx._borde(celda)

        assert self._lineas(celda) == ["lnL", "lnR", "lnT", "lnB"]

    def test_la_celda_absorbida_por_un_merge_tambien_lleva_borde(self, tabla):
        """(1,0) y (0,2) las come un merge y no pasan por `_pintar`.

        PowerPoint y LibreOffice no coinciden en de donde toman el borde
        exterior de un rango combinado, asi que se pintan igual."""
        assert set(self._lineas(tabla.cell(1, 0))) == {"lnL", "lnR", "lnT", "lnB"}
        assert set(self._lineas(tabla.cell(0, 2))) == {"lnL", "lnR", "lnT", "lnB"}
