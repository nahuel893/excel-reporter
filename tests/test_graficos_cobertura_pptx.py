"""Tests for pptx_builder — Marca.pptx and Generico.pptx decks."""
import pytest
from pptx import Presentation

from src.services.graficos_cobertura.pptx_builder import build_pptx, build_decks


def _make_png(path):
    """Create a minimal valid 1x1 PNG for testing."""
    from PIL import Image
    img = Image.new("RGB", (10, 10), "white")
    img.save(path, "PNG")


def _populate_png_dir(png_dir, zonas, genericos, include_comparacion=True):
    """Create a PNG file per (zona, generico) combination (and comparacion)."""
    png_dir.mkdir(parents=True, exist_ok=True)
    for zona_name, zona_slug in zonas:
        for gen_name, gen_slug in genericos:
            _make_png(png_dir / f"cobertura_{zona_slug}_{gen_slug}.png")
            if include_comparacion:
                _make_png(png_dir / f"comparacion_{zona_slug}_{gen_slug}.png")


class TestBuildPptx:
    """RF-016, RF-017: build_pptx creates a deck with expected slides."""

    _ZONAS = [
        ("NOA NORTE", "noa_norte"),
        ("SALTA CAPITAL", "salta_capital"),
    ]
    _GENERICOS = [
        ("CERVEZAS", "cervezas"),
    ]

    def test_creates_pptx_file(self, tmp_path):
        png_dir = tmp_path / "png"
        _populate_png_dir(png_dir, self._ZONAS, self._GENERICOS)

        out = build_pptx(
            genericos=self._GENERICOS,
            zonas=self._ZONAS,
            png_dir=png_dir,
            output_path=tmp_path / "Marca.pptx",
        )
        assert out.exists()
        assert out.suffix == ".pptx"

    def test_slide_count_matches_pngs(self, tmp_path):
        """With 2 zones × 1 generico and include_comparacion=True:
        2 cobertura slides + 2 comparacion slides = 4."""
        png_dir = tmp_path / "png"
        _populate_png_dir(png_dir, self._ZONAS, self._GENERICOS)

        out = build_pptx(
            genericos=self._GENERICOS,
            zonas=self._ZONAS,
            png_dir=png_dir,
            output_path=tmp_path / "deck.pptx",
        )
        prs = Presentation(str(out))
        assert len(prs.slides) == 4

    def test_include_comparacion_false_halves_slides(self, tmp_path):
        png_dir = tmp_path / "png"
        _populate_png_dir(png_dir, self._ZONAS, self._GENERICOS)

        out = build_pptx(
            genericos=self._GENERICOS,
            zonas=self._ZONAS,
            png_dir=png_dir,
            output_path=tmp_path / "deck_no_comp.pptx",
            include_comparacion=False,
        )
        prs = Presentation(str(out))
        # Only cobertura slides (no comparacion)
        assert len(prs.slides) == 2

    def test_missing_png_skipped_gracefully(self, tmp_path):
        png_dir = tmp_path / "png"
        png_dir.mkdir()
        # Create only ONE png of the two expected
        _make_png(png_dir / "cobertura_noa_norte_cervezas.png")

        out = build_pptx(
            genericos=self._GENERICOS,
            zonas=self._ZONAS,  # expects 2 zones
            png_dir=png_dir,
            output_path=tmp_path / "partial.pptx",
            include_comparacion=False,
        )
        prs = Presentation(str(out))
        # Only 1 slide because one png was missing
        assert len(prs.slides) == 1

    def test_slide_has_title_textbox(self, tmp_path):
        png_dir = tmp_path / "png"
        _populate_png_dir(png_dir, self._ZONAS, self._GENERICOS, include_comparacion=False)

        out = build_pptx(
            genericos=self._GENERICOS, zonas=self._ZONAS,
            png_dir=png_dir, output_path=tmp_path / "deck.pptx",
            include_comparacion=False,
        )
        prs = Presentation(str(out))
        slide = prs.slides[0]
        # Find at least one text frame containing the expected title substring
        titles = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                titles.append(shape.text_frame.text)
        combined = " | ".join(titles)
        assert "NOA NORTE" in combined
        assert "CERVEZAS" in combined


class TestBuildDecks:
    """RF-016, RF-017: build_decks produces Marca.pptx + Generico.pptx."""

    def test_creates_both_decks(self, tmp_path):
        # Create PNGs for Marca deck (cervezas + aguas) and Generico (all)
        png_dir = tmp_path / "png"
        zonas = [("NOA NORTE", "noa_norte")]
        all_gens = [
            ("CERVEZAS", "cervezas"),
            ("AGUAS SABORIZADAS", "aguas_saborizadas"),
            ("AGUAS MINERAL", "aguas_mineral"),
            ("SIDRAS Y LICORES", "sidras_y_licores"),
            ("VINOS CCU", "vinos_ccu"),
        ]
        _populate_png_dir(png_dir, zonas, all_gens)

        out_dir = tmp_path / "out"
        result = build_decks(png_dir=png_dir, output_dir=out_dir, con_aguas=True)

        assert "marca" in result
        assert "generico" in result
        assert result["marca"].exists()
        assert result["generico"].exists()

    def test_con_aguas_false_excludes_aguas_from_marca_deck(self, tmp_path):
        png_dir = tmp_path / "png"
        zonas = [("NOA NORTE", "noa_norte")]
        gens = [
            ("CERVEZAS", "cervezas"),
            ("AGUAS SABORIZADAS", "aguas_saborizadas"),
        ]
        _populate_png_dir(png_dir, zonas, gens, include_comparacion=False)

        out_dir = tmp_path / "out"
        result = build_decks(png_dir=png_dir, output_dir=out_dir, con_aguas=False)

        prs_marca = Presentation(str(result["marca"]))
        # With con_aguas=False, marca deck has only CERVEZAS (1 zone × 1 gen = 1 slide)
        assert len(prs_marca.slides) == 1
