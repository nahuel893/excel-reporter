"""Tests for pptx_builder — per-sucursal PPTX generation (T-005)."""
import pytest
from pathlib import Path
from pptx import Presentation

from src.services.graficos_cobertura.pptx_builder import build_deck_sucursal


def _make_png(path: Path) -> None:
    """Create a minimal valid 1x1 PNG for testing."""
    from PIL import Image
    img = Image.new("RGB", (10, 10), "white")
    img.save(path, "PNG")


def _populate_sucursal_pngs(
    png_dir: Path,
    zona_slug: str,
    gen_slug: str,
    sucursal_slug: str,
    include_comparacion: bool = True,
) -> dict[str, dict[str, Path]]:
    """Create zone + sucursal PNGs and return {scope: {type: path}}."""
    png_dir.mkdir(parents=True, exist_ok=True)
    zone = {}
    zone["cobertura"] = png_dir / f"cobertura_{zona_slug}_{gen_slug}.png"
    _make_png(zone["cobertura"])
    if include_comparacion:
        zone["comparacion"] = png_dir / f"comparacion_{zona_slug}_{gen_slug}.png"
        _make_png(zone["comparacion"])

    sucursal = {}
    sucursal["cobertura"] = png_dir / f"cobertura_{zona_slug}_suc{sucursal_slug}_{gen_slug}.png"
    _make_png(sucursal["cobertura"])
    if include_comparacion:
        sucursal["comparacion"] = png_dir / f"comparacion_{zona_slug}_suc{sucursal_slug}_{gen_slug}.png"
        _make_png(sucursal["comparacion"])

    return {"zone": zone, "sucursal": sucursal}


class TestBuildDeckSucursal:
    """T-005: build_deck_sucursal creates PPTX with zone + sucursal slides."""

    def test_creates_pptx_file(self, tmp_path: Path):
        pngs = _populate_sucursal_pngs(
            tmp_path / "png", "salta_capital", "cervezas", "6",
        )
        out = build_deck_sucursal(
            zona_slug="salta_capital",
            generico="CERVEZAS",
            sucursal_id=6,
            sucursal_nombre="SUCURSAL METAN",
            png_paths_zone=pngs["zone"],
            png_paths_sucursal=pngs["sucursal"],
            output_path=tmp_path / "cobertura_salta_capital_suc6.pptx",
        )
        assert out.exists()
        assert out.suffix == ".pptx"

    def test_slide_count_with_comparacion(self, tmp_path: Path):
        """With both cobertura + comparacion for zone and sucursal = 4 slides."""
        pngs = _populate_sucursal_pngs(
            tmp_path / "png", "salta_capital", "cervezas", "6",
        )
        out = build_deck_sucursal(
            zona_slug="salta_capital",
            generico="CERVEZAS",
            sucursal_id=6,
            sucursal_nombre="SUCURSAL METAN",
            png_paths_zone=pngs["zone"],
            png_paths_sucursal=pngs["sucursal"],
            output_path=tmp_path / "deck.pptx",
        )
        prs = Presentation(str(out))
        # Zone cobertura + Zone comparacion + Sucursal cobertura + Sucursal comparacion = 4
        assert len(prs.slides) == 4

    def test_slide_count_cobertura_only(self, tmp_path: Path):
        """Without comparacion PNGs = 2 slides (zone + sucursal cobertura)."""
        pngs = _populate_sucursal_pngs(
            tmp_path / "png", "salta_capital", "cervezas", "6",
            include_comparacion=False,
        )
        out = build_deck_sucursal(
            zona_slug="salta_capital",
            generico="CERVEZAS",
            sucursal_id=6,
            sucursal_nombre="SUCURSAL METAN",
            png_paths_zone=pngs["zone"],
            png_paths_sucursal=pngs["sucursal"],
            output_path=tmp_path / "deck_no_comp.pptx",
        )
        prs = Presentation(str(out))
        assert len(prs.slides) == 2

    def test_zone_slides_before_sucursal_slides(self, tmp_path: Path):
        """Zone overview slides come first, sucursal detail slides after."""
        pngs = _populate_sucursal_pngs(
            tmp_path / "png", "noa_norte", "cervezas", "1",
        )
        out = build_deck_sucursal(
            zona_slug="noa_norte",
            generico="CERVEZAS",
            sucursal_id=1,
            sucursal_nombre="CASA CENTRAL",
            png_paths_zone=pngs["zone"],
            png_paths_sucursal=pngs["sucursal"],
            output_path=tmp_path / "deck.pptx",
        )
        prs = Presentation(str(out))
        # Slide 1 = zone cobertura, Slide 2 = zone comparacion,
        # Slide 3 = sucursal cobertura, Slide 4 = sucursal comparacion
        titles = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    titles.append(shape.text_frame.text)

        # Zone slides should contain zona name but NOT sucursal name
        zone_titles = [t for t in titles if "NOA NORTE" in t]
        sucursal_titles = [t for t in titles if "CASA CENTRAL" in t]
        assert len(zone_titles) >= 2
        assert len(sucursal_titles) >= 2
        # Zone titles should appear before sucursal titles
        zone_idx = min(i for i, t in enumerate(titles) if "NOA NORTE" in t)
        suc_idx = min(i for i, t in enumerate(titles) if "CASA CENTRAL" in t)
        assert zone_idx < suc_idx

    def test_missing_png_skipped_gracefully(self, tmp_path: Path):
        """If some PNGs don't exist, only existing ones get slides."""
        png_dir = tmp_path / "png"
        png_dir.mkdir()
        # Only create zone cobertura + sucursal cobertura (no comparacion)
        zone = {"cobertura": png_dir / "cobertura_noa_norte_cervezas.png"}
        _make_png(zone["cobertura"])
        sucursal = {"cobertura": png_dir / "cobertura_noa_norte_suc1_cervezas.png"}
        _make_png(sucursal["cobertura"])

        out = build_deck_sucursal(
            zona_slug="noa_norte",
            generico="CERVEZAS",
            sucursal_id=1,
            sucursal_nombre="CASA CENTRAL",
            png_paths_zone=zone,
            png_paths_sucursal=sucursal,
            output_path=tmp_path / "partial.pptx",
        )
        prs = Presentation(str(out))
        # Only 2 slides (the 2 that have existing PNGs)
        assert len(prs.slides) == 2

    def test_output_in_sucursales_subdir(self, tmp_path: Path):
        """Output path can be placed in a sucursales/ subdirectory."""
        pngs = _populate_sucursal_pngs(
            tmp_path / "png", "salta_capital", "cervezas", "6",
        )
        out_dir = tmp_path / "sucursales"
        out_dir.mkdir()
        out = build_deck_sucursal(
            zona_slug="salta_capital",
            generico="CERVEZAS",
            sucursal_id=6,
            sucursal_nombre="SUCURSAL METAN",
            png_paths_zone=pngs["zone"],
            png_paths_sucursal=pngs["sucursal"],
            output_path=out_dir / "cobertura_salta_capital_suc6.pptx",
        )
        assert out.exists()
        assert out.parent.name == "sucursales"

    def test_different_generico(self, tmp_path: Path):
        """Works with non-CERVEZAS genericos like VINOS CCU."""
        pngs = _populate_sucursal_pngs(
            tmp_path / "png", "jujuy_interior", "vinos_ccu", "9",
        )
        out = build_deck_sucursal(
            zona_slug="jujuy_interior",
            generico="VINOS CCU",
            sucursal_id=9,
            sucursal_nombre="SUCURSAL JUJUY",
            png_paths_zone=pngs["zone"],
            png_paths_sucursal=pngs["sucursal"],
            output_path=tmp_path / "cobertura_jujuy_interior_suc9.pptx",
        )
        prs = Presentation(str(out))
        # Zone cobertura + comparacion + sucursal cobertura + comparacion = 4
        assert len(prs.slides) == 4
        titles = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    titles.append(shape.text_frame.text)
        combined = " | ".join(titles)
        assert "VINOS CCU" in combined
        assert "JUJUY INTERIOR" in combined
        assert "SUCURSAL JUJUY" in combined

    def test_all_missing_pngs_produces_empty_deck(self, tmp_path: Path):
        """If no PNGs exist, deck is created with 0 slides (no crash)."""
        zone = {
            "cobertura": tmp_path / "nonexistent_zone_cob.png",
            "comparacion": tmp_path / "nonexistent_zone_comp.png",
        }
        sucursal = {
            "cobertura": tmp_path / "nonexistent_suc_cob.png",
            "comparacion": tmp_path / "nonexistent_suc_comp.png",
        }
        out = build_deck_sucursal(
            zona_slug="noa_norte",
            generico="CERVEZAS",
            sucursal_id=1,
            sucursal_nombre="CASA CENTRAL",
            png_paths_zone=zone,
            png_paths_sucursal=sucursal,
            output_path=tmp_path / "empty.pptx",
        )
        assert out.exists()
        prs = Presentation(str(out))
        assert len(prs.slides) == 0

    def test_slide_dimensions_match_standard(self, tmp_path: Path):
        """Per-sucursal deck uses the same slide dimensions as standard decks."""
        pngs = _populate_sucursal_pngs(
            tmp_path / "png", "salta_capital", "cervezas", "6",
        )
        from src.services.graficos_cobertura.constants import (
            PPTX_SLIDE_WIDTH_IN,
            PPTX_SLIDE_HEIGHT_IN,
        )
        out = build_deck_sucursal(
            zona_slug="salta_capital",
            generico="CERVEZAS",
            sucursal_id=6,
            sucursal_nombre="SUCURSAL METAN",
            png_paths_zone=pngs["zone"],
            png_paths_sucursal=pngs["sucursal"],
            output_path=tmp_path / "dims.pptx",
        )
        prs = Presentation(str(out))
        from pptx.util import Inches
        assert prs.slide_width == Inches(PPTX_SLIDE_WIDTH_IN)
        assert prs.slide_height == Inches(PPTX_SLIDE_HEIGHT_IN)
